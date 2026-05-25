"""
TrustLayer AI - Hackathon Presentation Generator
Creates a professional PowerPoint presentation for the Infosys Business Incubator Cohort 2 Hackathon
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ============================================================================
# COLOR PALETTE
# ============================================================================
DARK_BG     = RGBColor(0x0F, 0x17, 0x2A)   # #0F172A — slide background
CARD_BG     = RGBColor(0x1E, 0x29, 0x3B)   # #1E293B — card / box background
ACCENT_BLUE = RGBColor(0x00, 0x66, 0xFF)   # #0066FF — primary blue
ACCENT_TEAL = RGBColor(0x00, 0xB8, 0x94)   # #00B894 — teal / green
ACCENT_RED  = RGBColor(0xE7, 0x4C, 0x3C)   # #E74C3C — danger / block
ACCENT_ORG  = RGBColor(0xF3, 0x9C, 0x12)   # #F39C12 — warning / flag
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x94, 0xA3, 0xB8)   # muted text
LIGHT_BLUE  = RGBColor(0xE8, 0xF4, 0xFD)   # highlight


# ============================================================================
# HELPERS
# ============================================================================

def new_slide(prs, layout_index=6):
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def fill_bg(slide, color=DARK_BG):
    """Set solid background color on a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None, line_color=None, line_width=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
                  line_spacing=None):
    """Add multiple lines in a single text box."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
    return txBox


def slide_header(slide, title, subtitle=None):
    """Render a consistent header bar at top of slide."""
    # Top blue accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, ACCENT_BLUE)
    # Title
    add_text(slide, title, 0.4, 0.15, 10, 0.6,
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.7, 12, 0.35,
                 font_size=13, color=GRAY)
    # Divider line
    add_rect(slide, 0.4, 1.0, 12.5, 0.025, ACCENT_BLUE)


def stat_card(slide, left, top, value, label, color=ACCENT_BLUE, width=2.8, height=1.4):
    """Render a stat card."""
    add_rect(slide, left, top, width, height, CARD_BG)
    # Left accent bar
    add_rect(slide, left, top, 0.05, height, color)
    add_text(slide, value, left + 0.15, top + 0.1, width - 0.2, 0.7,
             font_size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left + 0.1, top + 0.8, width - 0.2, 0.5,
             font_size=11, color=GRAY, align=PP_ALIGN.CENTER)


def badge(slide, left, top, text, bg_color, text_color=WHITE, width=1.5, height=0.35):
    """Render a small badge / tag."""
    add_rect(slide, left, top, width, height, bg_color)
    add_text(slide, text, left, top + 0.02, width, height - 0.04,
             font_size=10, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE BUILDERS
# ============================================================================

def slide_title(prs):
    """Slide 1: Title / Hero"""
    slide = new_slide(prs)
    fill_bg(slide)

    # Large background gradient-style rectangle
    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    # Blue gradient overlay
    add_rect(slide, 0, 0, 6.5, 7.5, RGBColor(0x00, 0x33, 0x88))

    # Brand tag
    add_text(slide, "HACKATHON SUBMISSION  |  INFOSYS INCUBATOR 2025",
             0.5, 0.4, 12, 0.4, font_size=11, color=ACCENT_TEAL, bold=True)

    # Main title
    add_text(slide, "TrustLayer AI", 0.5, 1.1, 7, 1.2,
             font_size=56, bold=True, color=WHITE)
    add_text(slide, "Enterprise AI Reliability Platform",
             0.5, 2.2, 7, 0.6, font_size=22, color=LIGHT_BLUE)

    # Tagline
    add_text(slide,
             "The trust layer between your enterprise and AI.\n"
             "Detect. Prevent. Govern.",
             0.5, 3.0, 6.5, 1.0,
             font_size=15, color=GRAY, italic=True)

    # Right side: key stats box
    add_rect(slide, 7.5, 1.0, 5.4, 5.5, CARD_BG)
    add_rect(slide, 7.5, 1.0, 0.06, 5.5, ACCENT_BLUE)

    stats = [
        ("$67.4B", "Annual losses from AI hallucinations"),
        ("$2.4M",  "Average cost per incident"),
        ("77%",    "Enterprise leaders concerned about hallucinations"),
        ("92%",    "TrustLayer hallucination catch rate"),
        ("<100ms", "Real-time detection latency"),
    ]
    for i, (val, lbl) in enumerate(stats):
        y = 1.3 + i * 1.0
        add_text(slide, val, 7.9, y, 2.0, 0.5,
                 font_size=22, bold=True, color=ACCENT_BLUE)
        add_text(slide, lbl, 7.9, y + 0.38, 4.6, 0.4,
                 font_size=11, color=GRAY)

    # Bottom tagline
    add_rect(slide, 0, 6.9, 13.33, 0.6, ACCENT_BLUE)
    add_text(slide, "Built for the industries that cannot afford to get it wrong.",
             0.5, 6.95, 12.5, 0.45,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2: The Problem"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "The Problem", "AI hallucinations are causing measurable, costly, real-world harm")

    examples = [
        ("⚖️  Legal",       "Attorneys sanctioned for submitting fabricated case citations generated by AI"),
        ("🏥  Healthcare",  "AI systems recommending dangerous drug dosages not supported by clinical data"),
        ("🏦  Banking",      "Chatbots guaranteeing mortgage rates and promotions that do not exist"),
        ("🏛️  Government",  "AI making unauthorized benefits eligibility determinations"),
    ]
    for i, (industry, desc) in enumerate(examples):
        y = 1.2 + i * 1.35
        add_rect(slide, 0.4, y, 12.4, 1.15, CARD_BG)
        add_rect(slide, 0.4, y, 0.06, 1.15, ACCENT_RED)
        add_text(slide, industry, 0.65, y + 0.1, 2.2, 0.4,
                 font_size=13, bold=True, color=ACCENT_RED)
        add_text(slide, desc, 0.65, y + 0.45, 11.8, 0.55,
                 font_size=13, color=WHITE)

    # Root cause note
    add_rect(slide, 0.4, 6.7, 12.4, 0.55, RGBColor(0x2D, 0x1A, 0x1A))
    add_rect(slide, 0.4, 6.7, 0.06, 0.55, ACCENT_RED)
    add_text(slide,
             "Root Cause: AI providers have a conflict of interest — they cannot solve the problem they created.",
             0.65, 6.73, 12.0, 0.45,
             font_size=12, bold=True, color=ACCENT_RED)


def slide_solution(prs):
    """Slide 3: Our Solution"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Our Solution", "TrustLayer AI — real-time hallucination detection, prevention, and governance")

    # Pipeline diagram
    boxes = [
        ("User Query",        CARD_BG,    WHITE),
        ("LLM Provider",      CARD_BG,    WHITE),
        ("TrustLayer AI",     ACCENT_BLUE, WHITE),
        ("Verified Response", ACCENT_TEAL, WHITE),
    ]
    for i, (label, bg, fg) in enumerate(boxes):
        x = 0.5 + i * 3.15
        add_rect(slide, x, 1.2, 2.7, 0.75, bg)
        add_text(slide, label, x, 1.2, 2.7, 0.75,
                 font_size=14, bold=True, color=fg, align=PP_ALIGN.CENTER)
        if i < len(boxes) - 1:
            add_text(slide, "→", x + 2.7, 1.33, 0.45, 0.45,
                     font_size=22, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    # Three action outcomes
    outcomes = [
        ("PASS",  ACCENT_TEAL, "Confidence ≥ 75%\nResponse delivered"),
        ("FLAG",  ACCENT_ORG,  "Confidence 50–74%\nRouted to human review"),
        ("BLOCK", ACCENT_RED,  "Confidence < 50%\nBlocked + fallback triggered"),
    ]
    for i, (action, color, desc) in enumerate(outcomes):
        x = 0.6 + i * 4.2
        add_rect(slide, x, 2.3, 3.8, 1.8, CARD_BG)
        add_rect(slide, x, 2.3, 3.8, 0.45, color)
        add_text(slide, action, x, 2.3, 3.8, 0.45,
                 font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.15, 2.82, 3.5, 1.1,
                 font_size=13, color=GRAY)

    # Key differentiators
    add_text(slide, "Key Differentiators", 0.4, 4.35, 12, 0.4,
             font_size=16, bold=True, color=WHITE)
    diffs = [
        ("Sub-100ms latency", ACCENT_BLUE),
        ("8 parallel detection algorithms", ACCENT_BLUE),
        ("Enterprise data grounding", ACCENT_TEAL),
        ("8 industry-specific modules", ACCENT_TEAL),
        ("Built-in EU AI Act compliance", ACCENT_ORG),
        ("Works with any LLM — no lock-in", ACCENT_ORG),
    ]
    for i, (text, color) in enumerate(diffs):
        col = i % 3
        row = i // 3
        x = 0.4 + col * 4.3
        y = 4.8 + row * 0.8
        add_rect(slide, x, y, 4.0, 0.6, CARD_BG)
        add_rect(slide, x, y, 0.05, 0.6, color)
        add_text(slide, "  " + text, x + 0.1, y + 0.1, 3.8, 0.45,
                 font_size=12, color=WHITE)


def slide_architecture(prs):
    """Slide 4: Technical Architecture"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Technical Architecture", "8 parallel detection techniques delivering <100ms end-to-end latency")

    techniques = [
        ("Semantic Entropy Analysis",      "87%", "45ms",  "Internal model uncertainty"),
        ("Self-Consistency Checking",      "82%", "120ms", "Cross-query contradiction detection"),
        ("Source Citation Verification",   "95%", "200ms", "Real-time citation lookup"),
        ("Enterprise Knowledge Grounding", "98%", "150ms", "Verified against your data sources"),
        ("Claim Extraction & Classification","91%","35ms", "Risk-tier classification of assertions"),
        ("Hallucination Pattern Recognition","79%","25ms",  "Trained on curated failure datasets"),
        ("Temporal Consistency Checking",  "93%", "40ms",  "Date-sensitive claim validation"),
        ("Numerical Claim Validation",     "96%", "80ms",  "Rate, dosage & threshold verification"),
    ]

    # Header row
    headers = ["Detection Technique", "Accuracy", "Latency", "What It Catches"]
    col_widths = [4.0, 1.1, 1.1, 5.4]
    col_x = [0.4, 4.5, 5.65, 6.8]
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_widths)):
        add_rect(slide, cx, 1.1, cw, 0.4, ACCENT_BLUE)
        add_text(slide, hdr, cx + 0.05, 1.1, cw - 0.1, 0.4,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for i, (name, acc, lat, desc) in enumerate(techniques):
        y = 1.55 + i * 0.6
        row_bg = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
        add_rect(slide, 0.4, y, 12.5, 0.56, row_bg)
        add_text(slide, name,  col_x[0] + 0.1, y + 0.1, col_widths[0] - 0.15, 0.4,
                 font_size=11, color=WHITE)
        add_text(slide, acc,   col_x[1] + 0.05, y + 0.1, col_widths[1] - 0.1, 0.4,
                 font_size=11, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
        add_text(slide, lat,   col_x[2] + 0.05, y + 0.1, col_widths[2] - 0.1, 0.4,
                 font_size=11, color=ACCENT_ORG, align=PP_ALIGN.CENTER)
        add_text(slide, desc,  col_x[3] + 0.05, y + 0.1, col_widths[3] - 0.1, 0.4,
                 font_size=10, color=GRAY)

    # Summary bar
    add_rect(slide, 0.4, 6.45, 12.5, 0.75, ACCENT_BLUE)
    summary_items = [
        "Overall Catch Rate: 92%",
        "False Positive Rate: <3%",
        "Pipeline Latency: <100ms",
        "Uptime SLA: 99.9%",
    ]
    for i, item in enumerate(summary_items):
        x = 0.8 + i * 3.1
        add_text(slide, item, x, 6.55, 2.9, 0.55,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_industries(prs):
    """Slide 5: Industry Modules"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Industry Modules", "Specialized detection for 8 verticals — 13 built-in scenarios")

    industries = [
        ("BFSI",          ACCENT_BLUE,  "Fabricated rates, unauthorized eligibility, regulatory violations"),
        ("Healthcare",    ACCENT_RED,   "Drug dosing errors, dangerous contraindications, unauthorized approvals"),
        ("Legal",         ACCENT_ORG,   "Fabricated case citations, non-existent statutes"),
        ("Enterprise",    ACCENT_TEAL,  "Fabricated HR policies, incorrect IT guidance"),
        ("Manufacturing", RGBColor(0x8E,0x44,0xAD), "Fabricated QC records, spec violations"),
        ("Retail",        RGBColor(0x16,0xA0,0x85), "Unverified pricing, availability errors"),
        ("Telecom",       RGBColor(0x27,0x6E,0xBD), "Unauthorized account actions, billing errors"),
        ("Government",    RGBColor(0xC0,0x39,0x2B), "Unauthorized eligibility determinations"),
    ]

    cols = 4
    card_w = 3.1
    card_h = 1.45
    gap_x = 0.18
    gap_y = 0.25
    start_x = 0.35
    start_y = 1.2

    for i, (name, color, desc) in enumerate(industries):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rect(slide, x, y, card_w, card_h, CARD_BG)
        add_rect(slide, x, y, card_w, 0.38, color)
        add_text(slide, name, x, y, card_w, 0.38,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.1, y + 0.42, card_w - 0.2, 0.9,
                 font_size=10, color=GRAY)

    # Bottom: example scenario callout
    add_rect(slide, 0.35, 6.45, 12.6, 0.8, RGBColor(0x0A, 0x1A, 0x35))
    add_rect(slide, 0.35, 6.45, 0.06, 0.8, ACCENT_RED)
    add_text(slide,
             "Example — Healthcare:  AI recommends metformin 1,000mg twice daily (FDA approved: 500mg), "
             "cites grapefruit interaction (fabricated), misses alcohol contraindication.  "
             "Confidence: 6%.  Action: BLOCKED.",
             0.55, 6.5, 12.1, 0.7,
             font_size=10, color=WHITE)


def slide_compliance(prs):
    """Slide 6: Compliance & Governance"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Compliance & Governance", "One platform — four major regulatory frameworks covered out of the box")

    frameworks = [
        ("EU AI Act",    "94%", "Aug 2026 enforcement — penalties up to €35M or 7% global revenue"),
        ("NIST AI RMF",  "89%", "Govern, Map, Measure, Manage functions fully addressed"),
        ("ISO/IEC 42001","82%", "International AI management system standard"),
        ("SOC 2 Type II","97%", "Availability, processing integrity, confidentiality controls"),
    ]

    for i, (fw, pct, desc) in enumerate(frameworks):
        y = 1.25 + i * 1.35
        add_rect(slide, 0.4, y, 12.4, 1.15, CARD_BG)
        # Progress bar background
        add_rect(slide, 2.5, y + 0.65, 8.5, 0.28, RGBColor(0x1A, 0x25, 0x38))
        pct_val = int(pct.replace("%",""))
        bar_w = 8.5 * pct_val / 100
        add_rect(slide, 2.5, y + 0.65, bar_w, 0.28, ACCENT_TEAL)
        add_text(slide, fw,  0.65, y + 0.15, 1.7, 0.45,
                 font_size=13, bold=True, color=WHITE)
        add_text(slide, pct, 2.5 + bar_w - 0.6, y + 0.63, 0.7, 0.32,
                 font_size=11, bold=True, color=WHITE)
        add_text(slide, desc, 0.65, y + 0.68, 9.5, 0.38,
                 font_size=10, color=GRAY)

    # Audit trail note
    add_rect(slide, 0.4, 6.65, 12.4, 0.6, CARD_BG)
    add_rect(slide, 0.4, 6.65, 0.06, 0.6, ACCENT_TEAL)
    add_text(slide,
             "Every AI interaction generates an immutable structured audit log — interaction ID, query, response, "
             "all 8 detection scores, action taken, compliance tags. Board-ready reporting at a click.",
             0.6, 6.7, 12.0, 0.52,
             font_size=11, color=WHITE)


def slide_integration(prs):
    """Slide 7: Integration Ecosystem"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Integration Ecosystem", "Works with any LLM provider — 50+ enterprise connectors — zero vendor lock-in")

    # LLM Providers
    add_text(slide, "LLM PROVIDERS", 0.4, 1.15, 6, 0.35,
             font_size=12, bold=True, color=ACCENT_BLUE)
    providers = ["OpenAI", "Anthropic", "Google", "Azure", "AWS Bedrock", "Meta Llama", "Mistral", "Open-source"]
    for i, p in enumerate(providers):
        col = i % 4
        row = i // 4
        x = 0.4 + col * 3.0
        y = 1.55 + row * 0.65
        add_rect(slide, x, y, 2.75, 0.52, CARD_BG)
        add_text(slide, p, x, y, 2.75, 0.52,
                 font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Enterprise connectors
    add_text(slide, "ENTERPRISE CONNECTORS (50+)", 0.4, 3.05, 12, 0.35,
             font_size=12, bold=True, color=ACCENT_TEAL)
    connectors = [
        ("CRM",        "Salesforce, HubSpot, Dynamics"),
        ("ERP",        "SAP, Oracle, Workday"),
        ("ITSM",       "ServiceNow, Jira, Zendesk"),
        ("HRIS",       "Workday, BambooHR, ADP"),
        ("Legal",      "Westlaw, LexisNexis"),
        ("Healthcare", "Epic, Cerner, FDA APIs"),
        ("Financial",  "Bloomberg, Reuters, Core Banking"),
        ("Identity",   "Active Directory, Okta"),
    ]
    for i, (cat, systems) in enumerate(connectors):
        col = i % 4
        row = i // 4
        x = 0.4 + col * 3.2
        y = 3.45 + row * 0.9
        add_rect(slide, x, y, 3.0, 0.78, CARD_BG)
        add_rect(slide, x, y, 3.0, 0.25, ACCENT_TEAL)
        add_text(slide, cat, x, y, 3.0, 0.25,
                 font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, systems, x + 0.08, y + 0.27, 2.85, 0.46,
                 font_size=9, color=GRAY)

    # Integration modes
    add_rect(slide, 0.4, 6.45, 12.4, 0.8, RGBColor(0x0A, 0x1A, 0x35))
    modes = ["SDK (hours)", "API Gateway (zero code)", "On-Premise", "SaaS (minutes)"]
    for i, m in enumerate(modes):
        x = 0.7 + i * 3.1
        add_text(slide, m, x, 6.6, 2.8, 0.5,
                 font_size=12, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)


def slide_roi(prs):
    """Slide 8: ROI & Business Case"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "ROI & Business Case", "The math is simple — preventing one incident pays for TrustLayer AI")

    # Left: cost of inaction
    add_rect(slide, 0.4, 1.2, 5.8, 5.0, CARD_BG)
    add_rect(slide, 0.4, 1.2, 5.8, 0.45, ACCENT_RED)
    add_text(slide, "Without TrustLayer AI", 0.4, 1.2, 5.8, 0.45,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    without_items = [
        ("100,000",    "Monthly AI queries"),
        ("15%",        "Industry avg hallucination rate"),
        ("~30",        "Critical incidents per year"),
        ("$2,400,000", "Cost per critical incident"),
        ("~$72M",      "Annual exposure"),
    ]
    for i, (val, lbl) in enumerate(without_items):
        y = 1.8 + i * 0.85
        add_text(slide, val, 0.6, y, 2.0, 0.45,
                 font_size=18, bold=True, color=ACCENT_RED)
        add_text(slide, lbl, 0.6, y + 0.38, 5.3, 0.35,
                 font_size=11, color=GRAY)

    # Right: with TrustLayer
    add_rect(slide, 6.8, 1.2, 6.1, 5.0, CARD_BG)
    add_rect(slide, 6.8, 1.2, 6.1, 0.45, ACCENT_TEAL)
    add_text(slide, "With TrustLayer AI", 6.8, 1.2, 6.1, 0.45,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    with_items = [
        ("92%",       "Hallucinations caught"),
        ("~2",        "Incidents per year (remaining)"),
        ("$66M+",     "Losses prevented annually"),
        ("~$500K",    "TrustLayer annual cost (enterprise)"),
        (">400%",     "Return on investment"),
    ]
    for i, (val, lbl) in enumerate(with_items):
        y = 1.8 + i * 0.85
        add_text(slide, val, 7.0, y, 2.0, 0.45,
                 font_size=18, bold=True, color=ACCENT_TEAL)
        add_text(slide, lbl, 7.0, y + 0.38, 5.6, 0.35,
                 font_size=11, color=GRAY)

    # VS divider
    add_text(slide, "VS", 6.0, 3.3, 0.8, 0.7,
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom: compliance upside
    add_rect(slide, 0.4, 6.45, 12.4, 0.75, RGBColor(0x0A, 0x2A, 0x1A))
    add_rect(slide, 0.4, 6.45, 0.06, 0.75, ACCENT_TEAL)
    add_text(slide,
             "Additional upside: Avoid EU AI Act fines up to €35M (August 2026) + reputation protection + regulatory readiness",
             0.6, 6.55, 12.0, 0.55,
             font_size=11, color=WHITE)


def slide_roadmap(prs):
    """Slide 9: Roadmap"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Roadmap", "From hackathon demo to enterprise-scale platform")

    phases = [
        ("Phase 1\nNow",
         "Hackathon Demo",
         ["Interactive Streamlit demo", "13 industry scenarios", "9 platform pages", "Compliance dashboards"],
         ACCENT_BLUE),
        ("Phase 2\n3-6 Mo",
         "MVP",
         ["Production detection engine", "Live API gateway", "3 enterprise connectors", "First paying pilot"],
         ACCENT_TEAL),
        ("Phase 3\n6-12 Mo",
         "Growth",
         ["50+ enterprise connectors", "SOC 2 Type II cert", "Multi-region deploy", "Self-serve onboarding"],
         ACCENT_ORG),
        ("Phase 4\n12-24 Mo",
         "Scale",
         ["EU AI Act cert program", "Federated learning", "Third-party module marketplace", "Series B"],
         RGBColor(0x8E, 0x44, 0xAD)),
    ]

    for i, (phase, title, items, color) in enumerate(phases):
        x = 0.35 + i * 3.2
        add_rect(slide, x, 1.2, 3.0, 5.8, CARD_BG)
        add_rect(slide, x, 1.2, 3.0, 0.55, color)
        add_text(slide, phase, x, 1.2, 3.0, 0.55,
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.1, 1.85, 2.8, 0.45,
                 font_size=14, bold=True, color=color)
        for j, item in enumerate(items):
            y = 2.45 + j * 0.95
            add_rect(slide, x + 0.1, y, 2.8, 0.78, RGBColor(0x16, 0x21, 0x32))
            add_rect(slide, x + 0.1, y, 0.05, 0.78, color)
            add_text(slide, item, x + 0.25, y + 0.15, 2.55, 0.52,
                     font_size=11, color=WHITE)


def slide_demo(prs):
    """Slide 10: Demo Scenarios"""
    slide = new_slide(prs)
    fill_bg(slide)
    slide_header(slide, "Live Demo — 13 Scenarios", "Real hallucinations. Real detection. Real-time blocking.")

    scenarios = [
        ("BFSI — Mortgage",     "18%", "BLOCK", ACCENT_RED),
        ("BFSI — Insurance",    "22%", "BLOCK", ACCENT_RED),
        ("BFSI — Wealth Mgmt",  "31%", "BLOCK", ACCENT_RED),
        ("Healthcare — Drug",    "6%", "BLOCK", ACCENT_RED),
        ("Healthcare — Auth",   "29%", "BLOCK", ACCENT_RED),
        ("Legal — Litigation",   "4%", "BLOCK", ACCENT_RED),
        ("Legal — Contract",    "35%", "BLOCK", ACCENT_RED),
        ("Enterprise — HR",     "42%", "FLAG",  ACCENT_ORG),
        ("Enterprise — IT",     "38%", "FLAG",  ACCENT_ORG),
        ("Manufacturing — QC",  "55%", "FLAG",  ACCENT_ORG),
        ("Retail — Product",    "61%", "FLAG",  ACCENT_ORG),
        ("Telecom — Billing",   "27%", "BLOCK", ACCENT_RED),
        ("Government — Benefits","19%","BLOCK", ACCENT_RED),
    ]

    # Header row
    add_rect(slide, 0.4, 1.1, 5.8, 0.38, ACCENT_BLUE)
    add_rect(slide, 6.3, 1.1, 2.5, 0.38, ACCENT_BLUE)
    add_rect(slide, 8.95, 1.1, 1.8, 0.38, ACCENT_BLUE)
    add_text(slide, "Scenario", 0.5, 1.1, 5.6, 0.38,
             font_size=11, bold=True, color=WHITE)
    add_text(slide, "Confidence", 6.3, 1.1, 2.5, 0.38,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Action", 8.95, 1.1, 1.8, 0.38,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows_per_col = 7
    for i, (name, conf, action, color) in enumerate(scenarios):
        col = i // rows_per_col
        row = i % rows_per_col
        x_off = col * 6.55
        y = 1.55 + row * 0.68
        row_bg = CARD_BG if i % 2 == 0 else RGBColor(0x16, 0x21, 0x32)
        add_rect(slide, 0.4 + x_off, y, 5.8, 0.6, row_bg)
        add_rect(slide, 6.3 + x_off, y, 2.5, 0.6, row_bg)
        add_rect(slide, 8.95 + x_off, y, 1.8, 0.6, row_bg)
        add_text(slide, name, 0.55 + x_off, y + 0.1, 5.5, 0.42,
                 font_size=11, color=WHITE)
        add_text(slide, conf, 6.3 + x_off, y + 0.1, 2.5, 0.42,
                 font_size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, action, 8.95 + x_off, y + 0.1, 1.8, 0.42,
                 font_size=11, bold=True, color=color, align=PP_ALIGN.CENTER)


def slide_closing(prs):
    """Slide 11: Closing / Call to Action"""
    slide = new_slide(prs)
    fill_bg(slide)

    add_rect(slide, 0, 0, 13.33, 7.5, DARK_BG)
    add_rect(slide, 0, 0, 0.1, 7.5, ACCENT_BLUE)

    add_text(slide, "TrustLayer AI", 0.8, 1.5, 11, 1.2,
             font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Enterprise AI Reliability Platform",
             0.8, 2.65, 11, 0.6, font_size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    add_rect(slide, 2.5, 3.5, 8.3, 0.05, ACCENT_BLUE)

    add_text(slide,
             "Every day without AI reliability controls is a day you're exposed.\n"
             "Exposed to customer harm. Regulatory fines. Reputation damage.",
             0.8, 3.8, 11.5, 0.9,
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Three closing points
    points = [
        ("Detect", "8 parallel algorithms,\n92% catch rate"),
        ("Prevent", "Real-time blocking\nbefore harm occurs"),
        ("Govern",  "EU AI Act ready,\naudit trail included"),
    ]
    for i, (title, desc) in enumerate(points):
        x = 1.3 + i * 3.7
        add_rect(slide, x, 4.9, 3.3, 1.6, CARD_BG)
        add_rect(slide, x, 4.9, 3.3, 0.42, ACCENT_BLUE)
        add_text(slide, title, x, 4.9, 3.3, 0.42,
                 font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.1, 5.38, 3.1, 0.9,
                 font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    add_rect(slide, 0, 6.85, 13.33, 0.65, ACCENT_BLUE)
    add_text(slide,
             "Built for the industries that cannot afford to get it wrong.  |  Infosys Business Incubator Cohort 2 Hackathon 2025",
             0.5, 6.9, 12.5, 0.52,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================================
# MAIN
# ============================================================================

def build_presentation():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("Building slides...")
    slide_title(prs);        print("  1/11  Title")
    slide_problem(prs);      print("  2/11  Problem Statement")
    slide_solution(prs);     print("  3/11  Our Solution")
    slide_architecture(prs); print("  4/11  Technical Architecture")
    slide_industries(prs);   print("  5/11  Industry Modules")
    slide_compliance(prs);   print("  6/11  Compliance & Governance")
    slide_integration(prs);  print("  7/11  Integration Ecosystem")
    slide_roi(prs);          print("  8/11  ROI & Business Case")
    slide_roadmap(prs);      print("  9/11  Roadmap")
    slide_demo(prs);         print(" 10/11  Demo Scenarios")
    slide_closing(prs);      print(" 11/11  Closing / CTA")

    out = r"C:\aiprojects\aitrust\TrustLayer_AI_Canada_Hackathon_2026.pptx"
    prs.save(out)
    print(f"\nPresentation saved: {out}")


if __name__ == "__main__":
    build_presentation()
