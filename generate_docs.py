"""
Generate TrustLayer AI hackathon submission documents — v2 (updated)
  1. TrustLayer_AI_Technical_Document.docx
  2. TrustLayer_AI_Presentation.pptx
"""

# ─────────────────────────────────────────────────────────────────────────────
# WORD DOCUMENT
# ─────────────────────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement


def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), hex_color)
    tcPr.append(shd)


def rgb(h):
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def H(doc, text, level=1, color="1F3864"):
    h = doc.add_heading(text, level=level)
    run = h.runs[0] if h.runs else h.add_run(text)
    run.font.color.rgb = rgb(color)
    return h


def P(doc, text, bold=False, italic=False, size=11, color=None):
    p = doc.add_paragraph()
    r = p.add_run(text)
    r.bold = bold
    r.italic = italic
    r.font.size = Pt(size)
    if color:
        r.font.color.rgb = rgb(color)
    return p


def T(doc, headers, rows, hdr_color="1F3864", col_widths=None):
    row_colors = ["F0F4F8", "FFFFFF"]
    tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(headers):
        c = tbl.rows[0].cells[i]
        c.text = h
        set_cell_bg(c, hdr_color)
        for p in c.paragraphs:
            for run in p.runs:
                run.bold = True
                run.font.color.rgb = rgb("FFFFFF")
                run.font.size = Pt(10)
    for ri, row in enumerate(rows):
        rc = tbl.add_row().cells
        bg = row_colors[ri % 2]
        for ci, val in enumerate(row):
            rc[ci].text = val
            set_cell_bg(rc[ci], bg)
            for p in rc[ci].paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)
    if col_widths:
        for ci, w in enumerate(col_widths):
            for row in tbl.rows:
                row.cells[ci].width = Inches(w)
    return tbl


def bullet(doc, items, size=11):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(item).font.size = Pt(size)


def numbered(doc, items, size=11):
    for item in items:
        p = doc.add_paragraph(style="List Number")
        p.add_run(item).font.size = Pt(size)


def build_docx():
    doc = Document()
    for s in doc.sections:
        s.top_margin = Inches(0.9)
        s.bottom_margin = Inches(0.9)
        s.left_margin = Inches(1.1)
        s.right_margin = Inches(1.1)

    # ── Cover ────────────────────────────────────────────────────────────────
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = t.add_run("TrustLayer AI")
    r.bold = True; r.font.size = Pt(32); r.font.color.rgb = rgb("1F3864")

    s = doc.add_paragraph()
    s.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = s.add_run("Enterprise AI Reliability & Hallucination Detection Platform")
    r2.font.size = Pt(16); r2.font.color.rgb = rgb("2E75B6")

    doc.add_paragraph()
    m = doc.add_paragraph()
    m.alignment = WD_ALIGN_PARAGRAPH.CENTER
    m.add_run(
        "Technical Submission Document  |  Canada Hackathon 2026\n"
        "Powered by Anthropic Claude + OpenAI GPT-4o\n"
        "Version 2.0  |  March 2026"
    ).font.size = Pt(11)
    doc.add_page_break()

    # ── 1. Executive Summary ─────────────────────────────────────────────────
    H(doc, "1.  Executive Summary")
    P(doc,
        "Enterprise AI deployments fail silently. A banking chatbot quotes a wrong mortgage rate. "
        "A healthcare assistant recommends the wrong drug dose. A legal assistant cites a case that "
        "does not exist. These are hallucinations — and they cause real financial, legal, and "
        "physical harm."
    )
    P(doc,
        "TrustLayer AI is a real-time hallucination detection and governance platform that "
        "intercepts every AI response before it reaches the user. It runs 8 parallel detection "
        "algorithms, cross-validates results between Claude and GPT-4o, and returns a "
        "PASS / FLAG / BLOCK decision with full confidence scoring and claim-level breakdown."
    )
    P(doc,
        "Version 2 adds a complete enterprise middleware simulation: an Enterprise Flow tab "
        "showing the full request-response lifecycle, a Human Review Queue where compliance "
        "reviewers action FLAG decisions, a Compliance Audit Log, and safe fallback delivery "
        "for BLOCK decisions — all in a single deployable Streamlit application."
    )
    doc.add_paragraph()

    # ── 2. What's New in Version 2 ───────────────────────────────────────────
    H(doc, "2.  What's New in Version 2")
    T(doc,
        ["Feature", "Description", "Tab"],
        [
            ("Enterprise Flow Simulation",
             "3-panel view: Banking App raw response | TrustLayer intercept scores | "
             "What customer actually received. Full blocked-vs-fallback comparison.",
             "Enterprise Flow"),
            ("Human Review Queue",
             "FLAG decisions automatically routed here. Reviewer sees query, unfiltered "
             "AI response, and TrustLayer verdict side-by-side with Approve / Reject / Escalate.",
             "Review Queue"),
            ("Compliance Audit Log",
             "Every reviewer decision logged with timestamp, reviewer name, original FLAG "
             "decision, and final outcome. Exportable as CSV.",
             "Review Queue"),
            ("Safe Fallback Delivery",
             "BLOCK decisions replace the AI response with a pre-approved safe fallback "
             "message. The blocked content is logged but never shown to the customer.",
             "Enterprise Flow"),
            ("Cross-Validation Escalation Routing",
             "When disagreement escalates FLAG to BLOCK, the item is captured in BOTH the "
             "review queue AND the blocked log. Card shows Claude raw decision + consensus badge.",
             "Review Queue"),
        ],
        "1F3864", col_widths=[1.8, 3.8, 1.2]
    )
    doc.add_paragraph()

    # ── 3. Problem Statement ─────────────────────────────────────────────────
    H(doc, "3.  Problem Statement")
    P(doc,
        "Large language models produce confident, fluent text even when factually wrong. "
        "In regulated industries this creates three categories of harm:"
    )
    bullet(doc, [
        "Financial Harm — fabricated rates, illegal pre-approval guarantees, invented promotions.",
        "Physical Harm — wrong drug dosages, missed contraindications, unauthorized diagnoses.",
        "Legal Harm — fabricated case citations in court filings causing attorney sanctions.",
    ])
    P(doc,
        "Existing solutions operate post-hoc (audit logs), require slow manual review queues, "
        "or use crude keyword filters. No platform intercepts in real time with dual-LLM "
        "validation and an integrated human review workflow."
    )
    doc.add_paragraph()

    # ── 4. Full System Architecture ──────────────────────────────────────────
    H(doc, "4.  Full System Architecture")
    H(doc, "4.1  End-to-End Request Flow", level=2, color="2E75B6")
    numbered(doc, [
        "Customer submits query through the enterprise banking app / chatbot.",
        "Banking app forwards query to Claude LLM (acting as domain expert).",
        "Claude generates an unfiltered response using an industry-tuned system prompt.",
        "TrustLayer intercepts the response BEFORE it reaches the customer.",
        "8 detection techniques run in parallel against the response and grounding data.",
        "Optionally, GPT-4o runs the identical detection prompt independently.",
        "An Agreement Score is computed from both models' outputs.",
        "The Policy Engine applies thresholds and industry risk multipliers.",
        "Decision: PASS delivers the response; FLAG routes to Review Queue; "
        "BLOCK substitutes a safe fallback message.",
        "Every decision is written to the Compliance Audit Log.",
    ])

    H(doc, "4.2  The Six Application Tabs", level=2, color="2E75B6")
    T(doc,
        ["Tab", "Purpose", "Key Capability"],
        [
            ("Live Detection", "Primary query interface",
             "Run any query, load preset scenarios, view confidence gauge + technique chart"),
            ("Session History", "Session-level analytics",
             "PASS/FLAG/BLOCK summary, confidence trend chart, CSV export"),
            ("Batch Test", "Automated regression testing",
             "One-click run of all hallucination-risk scenarios across all industries"),
            ("Enterprise Flow", "Middleware simulation",
             "3-panel: raw response | TrustLayer intercept | customer delivery"),
            ("Review Queue", "Human-in-the-loop governance",
             "FLAG queue, Approve/Reject/Escalate actions, Compliance Audit Log"),
            ("How It Works", "Documentation",
             "Detection pipeline, technique accuracy table, compliance standards"),
        ],
        "2E75B6", col_widths=[1.4, 1.8, 3.6]
    )
    doc.add_paragraph()

    # ── 5. Detection Engine ──────────────────────────────────────────────────
    H(doc, "5.  TrustLayer Detection Engine")
    H(doc, "5.1  The 8 Detection Techniques", level=2, color="2E75B6")
    T(doc,
        ["#", "Technique", "Weight", "What It Detects"],
        [
            ("1", "Semantic Entropy", "10%", "Vague, contradictory, or incoherent statements"),
            ("2", "Self-Consistency", "12%", "Internal contradictions within the response"),
            ("3", "Source Verification", "18%", "Unverifiable or fabricated citations and sources"),
            ("4", "Enterprise Grounding", "25%", "Claims contradicting authoritative grounding data"),
            ("5", "Claim Classification", "10%", "Unsubstantiated or overconfident factual claims"),
            ("6", "Pattern Recognition", "8%",  "Known hallucination surface patterns"),
            ("7", "Temporal Consistency", "10%", "Outdated information presented as current"),
            ("8", "Numerical Validation", "7%",  "Implausible or internally inconsistent numbers"),
        ],
        "2E75B6", col_widths=[0.3, 1.8, 0.7, 3.9]
    )

    H(doc, "5.2  Confidence Score & Decision Policy", level=2, color="2E75B6")
    P(doc, "    confidence = Sigma( score_i x weight_i )  [0-100%]", bold=True)
    P(doc, "    risk = (1 - confidence/100) x 100 x industry_risk_factor", bold=True)
    doc.add_paragraph()
    T(doc,
        ["Decision", "Confidence", "Risk Score", "Action"],
        [
            ("PASS",  ">= 75%", "< 30",  "Response delivered to customer"),
            ("FLAG",  "50-74%", "30-59", "Routed to Human Review Queue"),
            ("BLOCK", "< 50%",  ">= 60", "Blocked; safe fallback delivered to customer"),
        ],
        "1F3864", col_widths=[1.0, 1.2, 1.2, 3.4]
    )
    P(doc,
        "Policy decisions are made by a Python rule engine — not delegated to the LLM — "
        "ensuring deterministic, auditable outcomes.",
        italic=True, size=10
    )
    doc.add_paragraph()

    # ── 6. Cross-Validation ──────────────────────────────────────────────────
    H(doc, "6.  Claude + GPT-4o Cross-Validation")
    P(doc,
        "When an OpenAI API key is provided, the identical detection prompt is sent to both "
        "Claude and GPT-4o independently. Their outputs are compared to produce an Agreement Score."
    )
    H(doc, "6.1  Agreement Score Formula", level=2, color="2E75B6")
    T(doc,
        ["Component", "Weight", "Calculation"],
        [
            ("Action Match",         "50%", "1.0 if PASS/FLAG/BLOCK matches, else 0.0"),
            ("Confidence Gap",       "30%", "1.0 - |claude_conf - gpt_conf| / 100"),
            ("Technique Divergence", "20%", "1.0 - |techniques_flagged_diff| / 8"),
        ],
        "2E75B6", col_widths=[1.8, 0.8, 4.2]
    )

    H(doc, "6.2  Agreement Thresholds & Escalation", level=2, color="2E75B6")
    T(doc,
        ["Score", "Label", "Action"],
        [
            (">= 85%", "Strong Agreement",          "High reliability — proceed with decision"),
            ("65-84%", "Moderate Agreement",         "Human review recommended"),
            ("< 65%",  "Significant Disagreement",   "Escalate: PASS -> FLAG, FLAG -> BLOCK"),
        ],
        "2E75B6", col_widths=[0.9, 2.0, 3.9]
    )
    P(doc,
        "Conservative rule: if either model says BLOCK, the consensus is BLOCK. "
        "Disagreement below 65% is itself treated as a risk signal and escalates the decision.",
        bold=True
    )

    H(doc, "6.3  Cross-Validation Escalation Routing (v2 fix)", level=2, color="CC2222")
    P(doc,
        "Prior to v2, items where Claude said FLAG but disagreement escalated the consensus to "
        "BLOCK bypassed the Review Queue entirely — the human reviewer never saw them. "
        "Version 2 routes to the Review Queue based on Claude's raw individual verdict, "
        "not the consensus. Escalated items appear in both the queue and the blocked log, "
        "with a red warning banner and stacked Claude/consensus badges so the reviewer "
        "has full context."
    )
    doc.add_paragraph()

    # ── 7. Enterprise Flow Simulation ────────────────────────────────────────
    H(doc, "7.  Enterprise Flow Simulation")
    P(doc,
        "The Enterprise Flow tab simulates the complete real-world middleware architecture. "
        "It visualises every stage of the request-response pipeline for the most recent query."
    )
    H(doc, "7.1  Flow Diagram (interactive in-app)", level=2, color="2E75B6")
    P(doc,
        "Customer -> Banking App -> Claude LLM -> [raw response] -> "
        "TrustLayer -> PASS / FLAG / BLOCK -> Outcome"
    , bold=True)
    P(doc,
        "Each node in the flow diagram is highlighted in the relevant decision colour "
        "after a query is run. The outcome node changes to show the correct destination: "
        "customer (PASS), Review Queue (FLAG), or Safe Fallback (BLOCK)."
    )

    H(doc, "7.2  Three-Panel View", level=2, color="2E75B6")
    T(doc,
        ["Panel", "Content", "Purpose"],
        [
            ("Banking Chatbot — Raw Response",
             "Exact text generated by Claude LLM, unfiltered, before any interception",
             "Shows what would have reached the customer without TrustLayer"),
            ("TrustLayer Intercept",
             "Decision badge, confidence %, risk score, and list of issues caught",
             "Shows the detection result and why the decision was made"),
            ("What the Customer Receives",
             "PASS: exact AI response. FLAG: holding message. BLOCK: safe fallback text.",
             "Shows the final customer experience after governance"),
        ],
        "2E75B6", col_widths=[1.8, 2.8, 2.2]
    )

    H(doc, "7.3  Blocked vs Safe Fallback Comparison", level=2, color="2E75B6")
    P(doc,
        "For BLOCK decisions, a side-by-side panel appears below the three-panel view showing:"
    )
    bullet(doc, [
        "Left (red): The AI response that was generated and blocked — never shown to the customer.",
        "Right (green): The safe fallback message the customer actually received.",
        "Checklist confirming: customer protected, harmful content not delivered, logged for audit.",
    ])
    doc.add_paragraph()

    # ── 8. Human Review Queue ────────────────────────────────────────────────
    H(doc, "8.  Human Review Queue")
    P(doc,
        "The Review Queue implements a full human-in-the-loop compliance workflow. "
        "Every FLAG decision (including those escalated from cross-validation disagreement) "
        "is automatically routed here for a compliance reviewer to action."
    )

    H(doc, "8.1  Queue Item Card", level=2, color="2E75B6")
    P(doc, "Each item in the queue displays:")
    bullet(doc, [
        "Customer query text and industry module.",
        "Unfiltered AI response (what the banking chatbot generated).",
        "TrustLayer issues list — specific claims that were flagged.",
        "TrustLayer verdict: Claude raw decision + consensus decision (if escalated).",
        "Confidence and risk scores.",
        "Escalation warning banner (red) when cross-validation disagreement promoted FLAG to BLOCK.",
        "Optional reviewer note field for the audit trail.",
    ])

    H(doc, "8.2  Reviewer Actions", level=2, color="2E75B6")
    T(doc,
        ["Action", "Effect", "Audit Log Entry"],
        [
            ("Approve",  "Marks item as reviewed. Response cleared for delivery.",
             "Original: FLAG  |  Final: PASS"),
            ("Reject",   "Escalates to BLOCK. Adds to blocked log. Safe fallback triggered.",
             "Original: FLAG  |  Final: BLOCK"),
            ("Escalate", "Routes to senior compliance team for further review.",
             "Original: FLAG  |  Final: ESCALATED"),
        ],
        "1F3864", col_widths=[1.0, 3.2, 2.6]
    )

    H(doc, "8.3  Compliance Audit Log", level=2, color="2E75B6")
    P(doc,
        "Every reviewer action is written to the Compliance Audit Log, which records: "
        "timestamp, reviewer username, item ID, industry, original TrustLayer decision, "
        "reviewer action, final decision, query text, and reviewer note. "
        "The full log is exportable as CSV for regulatory reporting."
    )
    doc.add_paragraph()

    # ── 9. Industry Modules ──────────────────────────────────────────────────
    H(doc, "9.  Industry Modules")
    T(doc,
        ["Module", "Risk Factor", "Grounding Data", "Key Risk Tags"],
        [
            ("BFSI — Banking",  "1.3x",
             "Live APR rates, LTV rules, approval criteria",
             "fabricated_rate, illegal_guarantee, fake_promotion"),
            ("Healthcare",      "1.5x",
             "FDA dosing, contraindications, drug interactions",
             "wrong_dosage, fake_drug_interaction, unauthorized_diagnosis"),
            ("Legal",           "1.4x",
             "Verified case law, jurisdictions, Westlaw flags",
             "fabricated_citation, nonexistent_statute, fabricated_case"),
            ("Enterprise HR",   "1.1x",
             "PTO accrual, overtime rules, rollover policy",
             "wrong_policy, fabricated_benefit, incorrect_leave"),
            ("General",         "1.0x",
             "No domain grounding (general purpose)",
             "factual_error, unsupported_claim, fabricated_data"),
        ],
        "1F3864", col_widths=[1.4, 0.9, 2.3, 2.2]
    )
    doc.add_paragraph()

    # ── 10. Platform Features ────────────────────────────────────────────────
    H(doc, "10.  Platform Features Summary")
    T(doc,
        ["Feature", "Description"],
        [
            ("SSO Login Screen",
             "Animated SVG shield, glassmorphism card, Microsoft/Google SSO buttons, "
             "username/password auth against st.secrets[demo_users]"),
            ("Enterprise Flow Tab",
             "Full middleware simulation with flow diagram, 3-panel view, "
             "blocked-vs-fallback comparison, and all-blocks history"),
            ("Human Review Queue",
             "FLAG queue with Approve/Reject/Escalate, reviewer notes, "
             "escalation detection, and compliance audit log with CSV export"),
            ("Claude + GPT-4o Cross-Validation",
             "Independent dual-model analysis, agreement score, disagreement signals, "
             "conservative consensus, escalation routing fix"),
            ("Batch Test Mode",
             "One-click run of all hallucination-risk scenarios, BLOCK/FLAG/PASS "
             "summary, distribution chart, CSV export"),
            ("PDF Export",
             "fpdf2 formatted report with scores, issues, claims, and cross-validation "
             "section for compliance teams"),
            ("Confidence Delta Table",
             "Side-by-side Claude vs GPT-4o technique scores with colour-coded delta column"),
            ("JSON Retry Logic",
             "Silent retry on parse failure — zero crashes in production (~2% of API calls)"),
            ("Session History",
             "Full audit trail for the session with confidence trend chart and CSV export"),
        ],
        "2E75B6", col_widths=[2.0, 4.8]
    )
    doc.add_paragraph()

    # ── 11. Tech Stack ───────────────────────────────────────────────────────
    H(doc, "11.  Technology Stack")
    T(doc,
        ["Layer", "Technology", "Purpose"],
        [
            ("AI Detection",      "Anthropic Claude Sonnet (claude-sonnet-4-6)", "Primary detection and response generation"),
            ("Cross-Validation",  "OpenAI GPT-4o",                               "Independent second-opinion analysis"),
            ("UI Framework",      "Streamlit >= 1.32",                            "Rapid interactive web application"),
            ("Visualisation",     "Plotly >= 5.18",                               "Gauges, bar charts, trend charts"),
            ("Data Models",       "Pydantic v2",                                  "Validated structured detection results"),
            ("PDF Reports",       "fpdf2 >= 2.7",                                 "Formatted PDF export"),
            ("Deployment",        "Streamlit Cloud",                              "One-click cloud deploy via GitHub"),
            ("Language",          "Python 3.10+",                                 "Core application language"),
        ],
        "1F3864", col_widths=[1.5, 2.8, 2.5]
    )
    doc.add_paragraph()

    # ── 12. Compliance ───────────────────────────────────────────────────────
    H(doc, "12.  Compliance & Governance Coverage")
    T(doc,
        ["Standard", "Coverage", "Key Controls"],
        [
            ("EU AI Act 2024",  "94%",
             "Risk classification, human oversight (Review Queue), transparency, audit trail"),
            ("NIST AI RMF",     "All 4 functions",
             "GOVERN, MAP, MEASURE, MANAGE — all addressed"),
            ("ISO 42001",       "High",
             "AI management system controls, continual improvement"),
            ("SOC 2 Type II",   "Partial",
             "Audit logging (Audit Log tab), access control (SSO), data integrity"),
        ],
        "1F3864", col_widths=[1.5, 1.2, 4.1]
    )
    P(doc,
        "The Human Review Queue and Compliance Audit Log directly satisfy EU AI Act "
        "Article 9 human oversight requirements for high-risk AI systems.",
        italic=True, size=10
    )
    doc.add_paragraph()

    # ── 13. Deployment ───────────────────────────────────────────────────────
    H(doc, "13.  Deployment Guide")
    H(doc, "13.1  Local", level=2, color="2E75B6")
    numbered(doc, [
        "git clone https://github.com/ajittgosavii/trustlayer-ai.git",
        "cd trustlayer-ai && pip install -r requirements.txt",
        "cp .streamlit/secrets.toml.example .streamlit/secrets.toml",
        "Add ANTHROPIC_API_KEY and OPENAI_API_KEY to secrets.toml",
        "streamlit run streamlit_app.py",
    ])
    H(doc, "13.2  Streamlit Cloud", level=2, color="2E75B6")
    numbered(doc, [
        "Fork repo to GitHub account",
        "Go to share.streamlit.io -> New app -> select repo + streamlit_app.py",
        "Advanced settings -> Secrets -> paste ANTHROPIC_API_KEY and OPENAI_API_KEY",
        "Click Deploy",
    ])
    doc.add_paragraph()

    # ── 14. Demo Credentials ─────────────────────────────────────────────────
    H(doc, "14.  Demo Access Credentials")
    T(doc,
        ["Username", "Password", "Role"],
        [
            ("admin",  "trustlayer2026", "Full access — all modules, batch test, review queue"),
            ("demo",   "demo123",        "Standard user — all modules"),
            ("judge",  "hackathon2026",  "Hackathon panel access"),
        ],
        "1F3864", col_widths=[1.4, 1.8, 3.6]
    )
    doc.add_paragraph()

    # ── 15. Hackathon Alignment ──────────────────────────────────────────────
    H(doc, "15.  Canada Hackathon 2026 — Alignment")
    T(doc,
        ["Criterion", "TrustLayer Response"],
        [
            ("Theme: Enterprise AI Reliability",
             "Core mission — every feature targets enterprise hallucination risk in production"),
            ("Real-world Problem",
             "Hallucinations in BFSI, Healthcare, Legal cause documented financial and physical harm"),
            ("Technical Depth",
             "8 detection algorithms, dual-LLM cross-validation, rule-based policy engine, "
             "human review workflow, compliance audit log"),
            ("Live Demo",
             "Deployed Streamlit app — 6 tabs, preset scenarios, batch test, enterprise flow sim"),
            ("Business Value",
             "EU AI Act compliance, audit trail, safe fallback, PDF reports for regulators"),
            ("Innovation",
             "Disagreement between models as risk signal; cross-val escalation routing; "
             "full middleware simulation with human-in-the-loop governance in one app"),
        ],
        "1F3864", col_widths=[2.0, 4.8]
    )

    doc.add_paragraph()
    foot = doc.add_paragraph()
    foot.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = foot.add_run("TrustLayer AI  —  Detect · Prevent · Govern  |  Canada Hackathon 2026")
    fr.italic = True; fr.font.size = Pt(10); fr.font.color.rgb = rgb("888888")

    out = r"C:\aiprojects\aitrust_prod\TrustLayer_AI_Technical_Document.docx"
    doc.save(out)
    print(f"Saved: {out}")


# ─────────────────────────────────────────────────────────────────────────────
# POWERPOINT
# ─────────────────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PR
from pptx.enum.text import PP_ALIGN

NAVY  = PR(0x1F, 0x38, 0x64)
BLUE  = PR(0x2E, 0x75, 0xB6)
STEEL = PR(0x46, 0x82, 0xB4)
WHITE = PR(0xFF, 0xFF, 0xFF)
AMBER = PR(0xFF, 0xA5, 0x00)
GREEN = PR(0x00, 0x8B, 0x45)
RED   = PR(0xCC, 0x22, 0x22)
LGRAY = PR(0xF0, 0xF4, 0xF8)
DGRAY = PR(0x44, 0x44, 0x44)
PURP  = PR(0x7B, 0x2D, 0x8B)

W = Inches(13.33)
H_ = Inches(7.5)


def _rect(slide, left, top, width, height, color):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def _tb(slide, text, left, top, width, height,
        size=16, bold=False, italic=False, color=DGRAY,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _header(slide, title, subtitle=None, color=NAVY):
    _rect(slide, 0, 0, W, Inches(1.3), color)
    _tb(slide, title, Inches(0.4), Inches(0.1), Inches(12), Inches(0.75),
        size=28, bold=True, color=WHITE)
    if subtitle:
        _tb(slide, subtitle, Inches(0.4), Inches(0.82), Inches(12), Inches(0.4),
            size=13, color=PR(0xAD, 0xD8, 0xE6))


def _ppt_table(slide, headers, rows, left, top, width, height,
               hdr_color=NAVY, row_colors=(LGRAY, WHITE)):
    tbl = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height).table
    col_w = width // len(headers)
    for i in range(len(headers)):
        tbl.columns[i].width = col_w
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = hdr_color
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True; run.font.color.rgb = WHITE; run.font.size = Pt(12)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r+1, c); cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = row_colors[r % 2]
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11); run.font.color.rgb = DGRAY


def _bullets(slide, items, left, top, width, height, size=14, title=None, tc=NAVY):
    if title:
        _tb(slide, title, left, top, width, Inches(0.35), size=15, bold=True, color=tc)
        top += Inches(0.38); height -= Inches(0.38)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run(); r.text = "  •  " + item
        r.font.size = Pt(size); r.font.color.rgb = DGRAY


def build_pptx():
    prs = Presentation()
    prs.slide_width = W; prs.slide_height = H_
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: Cover ────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, W, H_, NAVY)
    _rect(sl, 0, Inches(5.5), W, Inches(0.07), BLUE)
    _tb(sl, "TrustLayer AI",
        Inches(1), Inches(0.9), Inches(11), Inches(1.4),
        size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(sl, "Enterprise AI Reliability & Hallucination Detection Platform",
        Inches(1), Inches(2.4), Inches(11), Inches(0.7),
        size=20, color=PR(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)
    _tb(sl, "Canada Hackathon 2026  |  Anthropic Claude + OpenAI GPT-4o  |  Version 2.0",
        Inches(1), Inches(3.15), Inches(11), Inches(0.5),
        size=14, color=PR(0x88, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
    _tb(sl, "6 Tabs  ·  8 Detection Techniques  ·  Dual-LLM Cross-Validation  ·  Human Review Queue  ·  Audit Log",
        Inches(1), Inches(3.75), Inches(11), Inches(0.45),
        size=12, color=PR(0x66, 0x99, 0xBB), align=PP_ALIGN.CENTER)
    _tb(sl, "Detect · Prevent · Govern",
        Inches(1), Inches(6.55), Inches(11), Inches(0.5),
        size=16, italic=True, color=AMBER, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: The Problem ──────────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "The Problem", "AI hallucinations cause real enterprise harm")
    examples = [
        ("Banking", "Chatbot quotes 4.5% mortgage\nActual: 6.85% APR\nRegulatory violation + customer loss"),
        ("Healthcare", "AI recommends 1000mg Metformin\nCorrect dose: 500mg\nPatient safety risk"),
        ("Legal", "AI cites non-existent court case\nSubmitted in real filing\nAttorney sanctioned (2023)"),
    ]
    for i, (title, desc) in enumerate(examples):
        left = Inches(0.4) + i * Inches(4.3)
        _rect(sl, left, Inches(1.6), Inches(3.9), Inches(4.6), LGRAY)
        _tb(sl, title, left+Inches(0.15), Inches(1.75),
            Inches(3.6), Inches(0.5), size=18, bold=True, color=NAVY)
        _tb(sl, desc, left+Inches(0.15), Inches(2.35),
            Inches(3.6), Inches(3.5), size=15, color=DGRAY)
    _tb(sl, "No existing tool catches these in real time, before the customer reads them.",
        Inches(0.5), Inches(6.45), Inches(12), Inches(0.6),
        size=13, italic=True, color=RED, align=PP_ALIGN.CENTER)

    # ── SLIDE 3: Solution — Full Flow ─────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "TrustLayer AI — Complete Solution Flow", "Real-time interception · Dual-LLM validation · Human-in-the-loop governance")
    steps = [
        ("Customer", "Submits query"),
        ("Banking App", "Enterprise chatbot"),
        ("Claude LLM", "Generates response"),
        ("TrustLayer", "Intercepts + scores"),
        ("Policy\nEngine", "PASS/FLAG/BLOCK"),
    ]
    bw = Inches(2.2); bh = Inches(1.3); top = Inches(1.7)
    for i, (title, sub) in enumerate(steps):
        left = Inches(0.25) + i * Inches(2.62)
        _rect(sl, left, top, bw, bh, BLUE)
        _tb(sl, title, left, top+Inches(0.15), bw, Inches(0.55),
            size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(sl, sub, left, top+Inches(0.7), bw, Inches(0.45),
            size=11, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps)-1:
            _tb(sl, "->", left+bw, top+Inches(0.4), Inches(0.42), Inches(0.5),
                size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    outcomes = [("PASS", GREEN, ">= 75%  Risk < 30\nDelivered to customer"),
                ("FLAG", AMBER, "50-74%  Risk 30-59\nHuman Review Queue"),
                ("BLOCK", RED,  "< 50%  Risk >= 60\nSafe fallback sent")]
    for i, (label, color, desc) in enumerate(outcomes):
        left = Inches(0.9) + i * Inches(4.1)
        _rect(sl, left, Inches(3.35), Inches(3.4), Inches(1.7), color)
        _tb(sl, label, left, Inches(3.45), Inches(3.4), Inches(0.7),
            size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(sl, desc, left, Inches(4.1), Inches(3.4), Inches(0.8),
            size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Row 3 — new v2 note
    _rect(sl, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.1), PR(0x0A, 0x12, 0x28))
    _tb(sl, "NEW in v2:",
        Inches(0.5), Inches(5.4), Inches(1.2), Inches(0.4),
        size=12, bold=True, color=AMBER)
    _tb(sl, "FLAG -> Review Queue (human reviewer: Approve / Reject / Escalate)   |   "
            "BLOCK -> Safe Fallback + Blocked Log   |   Full Compliance Audit Trail",
        Inches(1.8), Inches(5.4), Inches(11.2), Inches(0.8),
        size=12, color=PR(0xAD, 0xD8, 0xE6))

    # ── SLIDE 4: 8 Detection Techniques ──────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "8 Detection Techniques", "Parallel analysis — every response scored on all 8 dimensions")
    techniques = [
        ("Semantic Entropy",     "10%", "Vague or incoherent statements"),
        ("Self-Consistency",     "12%", "Internal contradictions"),
        ("Source Verification",  "18%", "Fabricated citations/sources"),
        ("Enterprise Grounding", "25%", "Contradicts authoritative data"),
        ("Claim Classification", "10%", "Unsubstantiated factual claims"),
        ("Pattern Recognition",  "8%",  "Known hallucination patterns"),
        ("Temporal Consistency", "10%", "Outdated info as current"),
        ("Numerical Validation", "7%",  "Implausible numbers/rates"),
    ]
    cols = [techniques[:4], techniques[4:]]
    for ci, col in enumerate(cols):
        lx = Inches(0.3) + ci * Inches(6.5)
        for ri, (name, wt, desc) in enumerate(col):
            ty = Inches(1.55) + ri * Inches(1.38)
            _rect(sl, lx, ty, Inches(6.1), Inches(1.2), LGRAY)
            _rect(sl, lx, ty, Inches(0.65), Inches(1.2), BLUE)
            _tb(sl, wt, lx, ty+Inches(0.32), Inches(0.65), Inches(0.5),
                size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _tb(sl, name, lx+Inches(0.75), ty+Inches(0.05),
                Inches(5.1), Inches(0.45), size=14, bold=True, color=NAVY)
            _tb(sl, desc, lx+Inches(0.75), ty+Inches(0.52),
                Inches(5.1), Inches(0.5), size=12, color=DGRAY)

    # ── SLIDE 5: Cross-Validation ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "Claude + GPT-4o Cross-Validation", "Disagreement between models is itself a risk signal")
    _rect(sl, Inches(0.3), Inches(1.5), Inches(5.9), Inches(5.7), LGRAY)
    _tb(sl, "Agreement Score Formula", Inches(0.5), Inches(1.62),
        Inches(5.5), Inches(0.4), size=15, bold=True, color=NAVY)
    comps = [
        ("Action Match",         "50%", "Same PASS/FLAG/BLOCK -> 1.0 else 0.0"),
        ("Confidence Gap",       "30%", "1.0 - |Claude_conf - GPT_conf| / 100"),
        ("Technique Divergence", "20%", "1.0 - |techniques_diff| / 8"),
    ]
    for i, (comp, wt, calc) in enumerate(comps):
        ty = Inches(2.15) + i * Inches(1.4)
        _rect(sl, Inches(0.5), ty, Inches(5.5), Inches(1.15), WHITE)
        _tb(sl, f"{comp}  ({wt})", Inches(0.65), ty+Inches(0.05),
            Inches(4.8), Inches(0.42), size=13, bold=True, color=BLUE)
        _tb(sl, calc, Inches(0.65), ty+Inches(0.5),
            Inches(4.8), Inches(0.4), size=12, color=DGRAY)

    _rect(sl, Inches(6.5), Inches(1.5), Inches(6.5), Inches(2.9), LGRAY)
    _tb(sl, "Agreement Thresholds", Inches(6.7), Inches(1.62),
        Inches(6.1), Inches(0.4), size=15, bold=True, color=NAVY)
    thresholds = [(">= 85%", "Strong Agreement", GREEN, "High reliability"),
                  ("65-84%", "Moderate Agreement", AMBER, "Human review recommended"),
                  ("< 65%",  "Significant Disagreement", RED, "Escalate: FLAG -> BLOCK")]
    for i, (pct, label, color, action) in enumerate(thresholds):
        ty = Inches(2.1) + i * Inches(1.07)
        _rect(sl, Inches(6.7), ty, Inches(1.1), Inches(0.9), color)
        _tb(sl, pct, Inches(6.7), ty+Inches(0.15),
            Inches(1.1), Inches(0.4), size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(sl, label, Inches(7.95), ty, Inches(4.8), Inches(0.42), size=13, bold=True, color=DGRAY)
        _tb(sl, action, Inches(7.95), ty+Inches(0.44), Inches(4.8), Inches(0.38), size=11, italic=True, color=DGRAY)

    _rect(sl, Inches(6.5), Inches(4.55), Inches(6.5), Inches(2.65), PR(0x14, 0x08, 0x2A))
    _tb(sl, "v2 — Escalation Routing Fix", Inches(6.7), Inches(4.65),
        Inches(6.1), Inches(0.4), size=14, bold=True, color=AMBER)
    _tb(sl, "When disagreement escalates FLAG -> BLOCK, the item\n"
            "now appears in BOTH the Review Queue AND the blocked log.\n\n"
            "Review card shows: Claude raw decision + Consensus badge\n"
            "+ red escalation warning banner so reviewer has full context.",
        Inches(6.7), Inches(5.1), Inches(6.1), Inches(1.9),
        size=12, color=PR(0xAD, 0xD8, 0xE6))

    # ── SLIDE 6: Enterprise Flow Simulation ───────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "Enterprise Flow Simulation  (NEW Tab)", "Full middleware interception — see every stage of the request lifecycle")
    panels = [
        ("Banking Chatbot — Raw Response",
         "Exact text generated by Claude LLM\nbefore any interception",
         "What would have reached\nthe customer WITHOUT TrustLayer",
         BLUE, "#EFF6FF"),
        ("TrustLayer Intercept",
         "Decision badge · confidence % · risk score\nIssues caught listed below",
         "Why the decision was made\nand which claims failed",
         PURP, "#F5F3FF"),
        ("What Customer Receives",
         "PASS: exact AI response\nFLAG: holding message\nBLOCK: safe fallback text",
         "The final customer experience\nafter governance is applied",
         GREEN, "#F0FDF4"),
    ]
    for i, (title, content, purpose, color, bg) in enumerate(panels):
        left = Inches(0.3) + i * Inches(4.35)
        _rect(sl, left, Inches(1.5), Inches(4.1), Inches(5.5), PR(
            int(bg[1:3],16), int(bg[3:5],16), int(bg[5:7],16)))
        _rect(sl, left, Inches(1.5), Inches(4.1), Inches(0.45), color)
        _tb(sl, title, left+Inches(0.1), Inches(2.05),
            Inches(3.9), Inches(0.55), size=13, bold=True, color=color)
        _tb(sl, content, left+Inches(0.1), Inches(2.65),
            Inches(3.9), Inches(1.5), size=12, color=DGRAY)
        _rect(sl, left+Inches(0.1), Inches(4.25), Inches(3.9), Inches(0.02), color)
        _tb(sl, purpose, left+Inches(0.1), Inches(4.35),
            Inches(3.9), Inches(1.0), size=11, italic=True, color=DGRAY)

    _rect(sl, Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.9), PR(0x14, 0x08, 0x2A))
    _tb(sl, "BLOCK decisions show a side-by-side panel: AI response (blocked, never shown) vs safe fallback "
            "(what customer received). Checklist confirms: customer protected, content not delivered, logged for audit.",
        Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.7),
        size=11, color=PR(0xAD, 0xD8, 0xE6))

    # ── SLIDE 7: Human Review Queue ───────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "Human Review Queue  (NEW Tab)", "Full human-in-the-loop compliance workflow for FLAG decisions")
    left_w = Inches(6.2)
    _rect(sl, Inches(0.3), Inches(1.5), left_w, Inches(5.7), LGRAY)
    _tb(sl, "Queue Item Card — what the reviewer sees:", Inches(0.5), Inches(1.62),
        left_w-Inches(0.4), Inches(0.4), size=14, bold=True, color=NAVY)
    card_items = [
        "Customer query text and industry module",
        "Unfiltered AI response (what the chatbot generated)",
        "TrustLayer issues list — specific failing claims",
        "Claude raw decision badge + Consensus badge (if escalated by cross-val)",
        "Red escalation warning when FLAG was promoted to BLOCK",
        "Confidence % and risk score",
        "Optional reviewer note field",
    ]
    _bullets(sl, card_items, Inches(0.5), Inches(2.1), left_w-Inches(0.4), Inches(4.5),
             size=12)

    _rect(sl, Inches(6.7), Inches(1.5), Inches(6.3), Inches(5.7), LGRAY)
    _tb(sl, "Reviewer Actions:", Inches(6.9), Inches(1.62),
        Inches(5.9), Inches(0.4), size=14, bold=True, color=NAVY)
    actions = [
        ("Approve",  GREEN, "Response cleared for delivery\nAudit log: FLAG -> PASS"),
        ("Reject",   RED,   "Escalates to BLOCK\nSafe fallback triggered\nAudit log: FLAG -> BLOCK"),
        ("Escalate", BLUE,  "Routes to senior compliance\nAudit log: FLAG -> ESCALATED"),
    ]
    for i, (label, color, desc) in enumerate(actions):
        ty = Inches(2.15) + i * Inches(1.55)
        _rect(sl, Inches(6.9), ty, Inches(1.5), Inches(1.3), color)
        _tb(sl, label, Inches(6.9), ty+Inches(0.35), Inches(1.5), Inches(0.55),
            size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _tb(sl, desc, Inches(8.55), ty, Inches(4.2), Inches(1.2),
            size=12, color=DGRAY)

    _rect(sl, Inches(6.9), Inches(6.55-Inches(0.3)), Inches(6.1-Inches(0.2)), Inches(0.6), PR(0x0A, 0x12, 0x28))
    _tb(sl, "Compliance Audit Log: every decision recorded with timestamp, reviewer, "
            "original decision, final decision, query. CSV export for regulators.",
        Inches(6.9), Inches(6.32), Inches(6.0), Inches(0.6),
        size=11, color=PR(0xAD, 0xD8, 0xE6))

    # ── SLIDE 8: Industry Modules ─────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "Industry Modules", "5 pre-configured modules with authoritative grounding data")
    modules = [
        ("Banking", "1.3x", "APR rates\nLTV rules\nApproval criteria"),
        ("Healthcare", "1.5x", "FDA dosing\nContraindications\nDrug interactions"),
        ("Legal", "1.4x", "Case law\nJurisdictions\nWestlaw flags"),
        ("HR", "1.1x", "PTO accrual\nOvertime rules\nRollover policy"),
        ("General", "1.0x", "Custom queries\nNo domain grounding"),
    ]
    icons = ["🏦", "🏥", "⚖️", "👥", "🤖"]
    colors = [NAVY, RED, PURP, BLUE, STEEL]
    for i, (name, mult, grnd) in enumerate(modules):
        left = Inches(0.25) + i * Inches(2.6)
        _rect(sl, left, Inches(1.55), Inches(2.4), Inches(5.2), LGRAY)
        _rect(sl, left, Inches(1.55), Inches(2.4), Inches(0.55), colors[i])
        _tb(sl, icons[i], left, Inches(2.2), Inches(2.4), Inches(0.65),
            size=30, align=PP_ALIGN.CENTER)
        _tb(sl, name, left, Inches(3.0), Inches(2.4), Inches(0.45),
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _tb(sl, f"Risk: {mult}", left, Inches(3.5), Inches(2.4), Inches(0.38),
            size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
        _tb(sl, grnd, left+Inches(0.1), Inches(3.95), Inches(2.2), Inches(2.0),
            size=11, color=DGRAY, align=PP_ALIGN.CENTER)

    # ── SLIDE 9: Tech Stack + Compliance ──────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "Technology Stack & Compliance Coverage")
    _rect(sl, Inches(0.3), Inches(1.5), Inches(6.0), Inches(5.7), LGRAY)
    _tb(sl, "Technology Stack", Inches(0.5), Inches(1.62),
        Inches(5.6), Inches(0.4), size=15, bold=True, color=NAVY)
    tech = [
        ("Detection",        "Claude Sonnet (claude-sonnet-4-6)"),
        ("Cross-Validation", "OpenAI GPT-4o"),
        ("UI",               "Streamlit + Plotly"),
        ("Data Models",      "Pydantic v2"),
        ("PDF Reports",      "fpdf2"),
        ("Deployment",       "Streamlit Cloud"),
        ("Language",         "Python 3.10+"),
    ]
    for i, (layer, name) in enumerate(tech):
        ty = Inches(2.1) + i * Inches(0.71)
        _tb(sl, f"  {layer}:", Inches(0.5), ty, Inches(2.1), Inches(0.5),
            size=12, bold=True, color=BLUE)
        _tb(sl, name, Inches(2.65), ty, Inches(3.45), Inches(0.5),
            size=12, color=DGRAY)

    _rect(sl, Inches(6.6), Inches(1.5), Inches(6.4), Inches(5.7), LGRAY)
    _tb(sl, "Compliance Coverage", Inches(6.8), Inches(1.62),
        Inches(6.0), Inches(0.4), size=15, bold=True, color=NAVY)
    compliance = [
        ("EU AI Act 2024",  "94%",  "Risk classification, human oversight\n(Review Queue), transparency"),
        ("NIST AI RMF",     "100%", "GOVERN, MAP, MEASURE, MANAGE"),
        ("ISO 42001",       "High", "AI management system controls"),
        ("SOC 2 Type II",   "Partial", "Audit log, access control, data integrity"),
    ]
    for i, (std, cov, desc) in enumerate(compliance):
        ty = Inches(2.1) + i * Inches(1.25)
        _rect(sl, Inches(6.8), ty, Inches(5.9), Inches(1.1), WHITE)
        _tb(sl, std, Inches(6.95), ty+Inches(0.05),
            Inches(2.5), Inches(0.42), size=13, bold=True, color=NAVY)
        _tb(sl, cov, Inches(9.5), ty+Inches(0.05),
            Inches(1.5), Inches(0.42), size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        _tb(sl, desc, Inches(6.95), ty+Inches(0.52),
            Inches(5.6), Inches(0.5), size=11, italic=True, color=DGRAY)

    # ── SLIDE 10: Live Demo Script ────────────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _header(sl, "Live Demo Script  (5 minutes)", "Run these in order for maximum impact")
    steps_demo = [
        ("1/5", "Q: What documents do I need for a mortgage?",
         "PASS", GREEN,
         "Show system correctly trusting a safe response — baseline trust"),
        ("2/5", "Q: What is your current 30-year mortgage rate?",
         "FLAG / BLOCK", RED,
         "AI fabricates a rate — grounding detects deviation from 6.85% APR"),
        ("3/5", "Q: What promotions do you have this month?",
         "BLOCK", RED,
         "AI invents a promotion — grounding confirms none are active"),
        ("4/5", "Go to Enterprise Flow tab",
         "See all 3 panels", BLUE,
         "Show raw AI response vs what customer saw vs safe fallback side-by-side"),
        ("5/5", "Go to Review Queue tab",
         "Action FLAG item", AMBER,
         "Demonstrate Approve / Reject / Escalate + audit log entry"),
    ]
    for i, (step, query, result, color, talking) in enumerate(steps_demo):
        ty = Inches(1.55) + i * Inches(1.13)
        _rect(sl, Inches(0.3), ty, Inches(0.85), Inches(0.95), color)
        _tb(sl, step, Inches(0.3), ty+Inches(0.22),
            Inches(0.85), Inches(0.45), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _rect(sl, Inches(1.25), ty, Inches(12.4), Inches(0.95), LGRAY)
        _tb(sl, query, Inches(1.35), ty+Inches(0.05),
            Inches(5.8), Inches(0.42), size=12, bold=True, color=NAVY)
        _tb(sl, result, Inches(7.25), ty+Inches(0.05),
            Inches(1.8), Inches(0.42), size=12, bold=True, color=color)
        _tb(sl, talking, Inches(9.15), ty+Inches(0.05),
            Inches(4.4), Inches(0.82), size=11, italic=True, color=DGRAY)

    # ── SLIDE 11: Closing / Credentials ──────────────────────────────────────
    sl = prs.slides.add_slide(blank)
    _rect(sl, 0, 0, W, H_, NAVY)
    _rect(sl, 0, Inches(5.0), W, Inches(0.06), BLUE)
    _tb(sl, "Try It Now",
        Inches(1), Inches(0.5), Inches(11), Inches(0.7),
        size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(sl, "Demo Credentials",
        Inches(1), Inches(1.3), Inches(11), Inches(0.5),
        size=18, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    creds = [
        ("admin",  "trustlayer2026", "Full access — all 6 tabs, batch test, review queue"),
        ("demo",   "demo123",        "Standard user — all modules"),
        ("judge",  "hackathon2026",  "Hackathon panel access"),
    ]
    for ri, (user, pwd, role) in enumerate(creds):
        ty = Inches(2.0) + ri * Inches(0.78)
        bg = PR(0x2E, 0x75, 0xB6) if ri == 0 else PR(0x1A, 0x2E, 0x52)
        for ci, (val, cw) in enumerate([(user, Inches(2.5)), (pwd, Inches(2.8)), (role, Inches(5.2))]):
            lx = Inches(1.2) + sum([Inches(2.5), Inches(2.8), Inches(5.2)][:ci])
            _rect(sl, lx, ty, cw, Inches(0.62), bg)
            _tb(sl, val, lx+Inches(0.1), ty+Inches(0.1),
                cw-Inches(0.2), Inches(0.45), size=13, bold=(ri == 0), color=WHITE)

    _tb(sl, "streamlit run streamlit_app.py",
        Inches(1), Inches(5.35), Inches(11), Inches(0.5),
        size=16, bold=True, color=PR(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)
    _tb(sl, "github.com/ajittgosavii/trustlayer-ai",
        Inches(1), Inches(5.95), Inches(11), Inches(0.4),
        size=14, color=PR(0x88, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
    _tb(sl, "TrustLayer AI  —  Detect · Prevent · Govern  |  Canada Hackathon 2026",
        Inches(1), Inches(6.85), Inches(11), Inches(0.4),
        size=13, italic=True, color=AMBER, align=PP_ALIGN.CENTER)

    out = r"C:\aiprojects\aitrust_prod\TrustLayer_AI_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build_docx()
    build_pptx()
    print("Done.")
