"""
TrustLayer AI - Detailed Hackathon Presentation Generator (v2)
~24 slides with deep content per section
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================================
# PALETTE
# ============================================================================
BG        = RGBColor(0x0F, 0x17, 0x2A)
CARD      = RGBColor(0x1E, 0x29, 0x3B)
CARD2     = RGBColor(0x16, 0x21, 0x32)
BLUE      = RGBColor(0x00, 0x66, 0xFF)
DARK_BLUE = RGBColor(0x00, 0x33, 0x88)
TEAL      = RGBColor(0x00, 0xB8, 0x94)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE    = RGBColor(0xF3, 0x9C, 0x12)
PURPLE    = RGBColor(0x8E, 0x44, 0xAD)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
LGRAY     = RGBColor(0xCB, 0xD5, 0xE1)
LBLUE     = RGBColor(0xE8, 0xF4, 0xFD)
GREEN     = RGBColor(0x27, 0xAE, 0x60)

# ============================================================================
# CORE HELPERS
# ============================================================================

def slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def rect(s, l, t, w, h, fill, line=None, lw=None):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw or 1)
    else: sh.line.fill.background()
    return sh

def txt(s, text, l, t, w, h, size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return tb

def mtxt(s, lines, l, t, w, h, size=12, bold=False, color=WHITE,
         align=PP_ALIGN.LEFT, spacing=None, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing: p.line_spacing = Pt(spacing)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Segoe UI"
    return tb

def hdr(s, title, subtitle=None):
    rect(s, 0, 0, 13.33, 0.07, BLUE)
    txt(s, title, 0.35, 0.12, 10, 0.6, size=26, bold=True)
    if subtitle:
        txt(s, subtitle, 0.35, 0.65, 12.5, 0.32, size=12, color=GRAY)
    rect(s, 0.35, 0.95, 12.6, 0.02, BLUE)

def card(s, l, t, w, h, fill=CARD):
    rect(s, l, t, w, h, fill)

def accent_card(s, l, t, w, h, color=BLUE, fill=CARD):
    rect(s, l, t, w, h, fill)
    rect(s, l, t, 0.05, h, color)

def tag(s, l, t, text, color=BLUE, w=1.6, h=0.3):
    rect(s, l, t, w, h, color)
    txt(s, text, l, t+0.02, w, h-0.04, size=9, bold=True, align=PP_ALIGN.CENTER)

def divider(s, y=7.2):
    rect(s, 0, y, 13.33, 0.3, DARK_BLUE)

# ============================================================================
# SLIDE 1 — TITLE / HERO
# ============================================================================
def s01_title(prs):
    s = slide(prs)
    rect(s, 0, 0, 7.0, 7.5, RGBColor(0x00, 0x1A, 0x55))
    rect(s, 0, 0, 0.12, 7.5, BLUE)

    txt(s, "Infosys Business Incubator Cohort 2", 0.4, 0.35, 8, 0.35,
        size=10, color=TEAL, bold=True)
    txt(s, "TrustLayer AI", 0.4, 0.85, 6.5, 1.3, size=54, bold=True)
    txt(s, "Enterprise AI Reliability Platform", 0.4, 2.1, 6.5, 0.6,
        size=20, color=LBLUE)
    rect(s, 0.4, 2.85, 5.5, 0.03, BLUE)
    txt(s, "Detect  ·  Prevent  ·  Govern", 0.4, 3.0, 6.0, 0.5,
        size=16, color=GRAY, italic=True)

    mtxt(s, [
        "The trust layer between your enterprise and AI.",
        "Real-time hallucination detection across 8 industries,",
        "8 LLM providers, and 4 compliance frameworks.",
    ], 0.4, 3.65, 6.4, 1.2, size=13, color=LGRAY)

    # Right panel
    rect(s, 7.4, 0.5, 5.6, 6.7, CARD)
    rect(s, 7.4, 0.5, 5.6, 0.45, BLUE)
    txt(s, "The Numbers", 7.4, 0.5, 5.6, 0.45, size=14, bold=True, align=PP_ALIGN.CENTER)

    stats = [
        ("$67.4B",   "Annual enterprise losses from AI hallucinations", RED),
        ("$2.4M",    "Average cost per critical incident",              ORANGE),
        ("77%",      "Enterprise leaders concerned about hallucinations", ORANGE),
        ("92%",      "TrustLayer hallucination catch rate",              TEAL),
        ("<100ms",   "Real-time detection pipeline latency",             TEAL),
        (">400%",    "Typical return on investment year 1",              GREEN),
    ]
    for i, (val, lbl, col) in enumerate(stats):
        y = 1.15 + i * 0.98
        rect(s, 7.6, y, 5.2, 0.82, CARD2)
        rect(s, 7.6, y, 0.05, 0.82, col)
        txt(s, val, 7.75, y + 0.06, 1.8, 0.4, size=22, bold=True, color=col)
        txt(s, lbl, 7.75, y + 0.44, 4.8, 0.32, size=10, color=GRAY)

    rect(s, 0, 7.2, 13.33, 0.3, BLUE)
    txt(s, "Built for the industries that cannot afford to get it wrong.",
        0.4, 7.22, 12.5, 0.26, size=12, bold=True, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2 — THE AI HALLUCINATION CRISIS (SCALE)
# ============================================================================
def s02_crisis_scale(prs):
    s = slide(prs)
    hdr(s, "The AI Hallucination Crisis", "Scale of the problem — why this cannot be ignored")

    # Big stat
    txt(s, "$67,400,000,000", 0.35, 1.1, 12.6, 1.0, size=44, bold=True,
        color=RED, align=PP_ALIGN.CENTER)
    txt(s, "Lost annually by enterprises due to AI hallucinations",
        0.35, 2.0, 12.6, 0.4, size=15, color=GRAY, align=PP_ALIGN.CENTER)

    rect(s, 0.35, 2.55, 12.6, 0.02, CARD)

    stats4 = [
        ("77%",    "of enterprise leaders are\nconcerned about hallucinations",    ORANGE),
        ("$2.4M",  "average cost of a single\ncritical hallucination incident",    RED),
        ("15%",    "industry average AI\nhallucination rate",                      ORANGE),
        ("~30",    "critical incidents per year for a\ntypical enterprise at scale",RED),
    ]
    for i, (val, lbl, col) in enumerate(stats4):
        x = 0.35 + i * 3.2
        rect(s, x, 2.75, 3.0, 1.7, CARD)
        rect(s, x, 2.75, 3.0, 0.06, col)
        txt(s, val, x, 2.9, 3.0, 0.7, size=34, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, lbl, x + 0.1, 3.62, 2.8, 0.7, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Timeline of incidents
    txt(s, "HIGH-PROFILE INCIDENTS TIMELINE", 0.35, 4.65, 12.6, 0.35,
        size=11, bold=True, color=BLUE)
    events = [
        ("2023",  "Mata v. Avianca — Attorneys sanctioned for AI-fabricated case citations"),
        ("2024",  "Healthcare AI tools flagged for dangerous dosage recommendations"),
        ("2024",  "Banking chatbots guaranteeing mortgage terms that did not exist"),
        ("2025",  "EU AI Act passes — August 2026 enforcement deadline announced"),
        ("2025",  "77% of enterprise leaders cite hallucinations as top AI risk"),
    ]
    for i, (year, event) in enumerate(events):
        y = 5.05 + i * 0.42
        rect(s, 0.35, y, 12.6, 0.37, CARD2 if i % 2 == 0 else CARD)
        tag(s, 0.4, y + 0.04, year, BLUE, w=0.7, h=0.28)
        txt(s, event, 1.2, y + 0.05, 11.5, 0.28, size=11, color=WHITE)

# ============================================================================
# SLIDE 3 — REAL-WORLD HALLUCINATION EXAMPLES
# ============================================================================
def s03_incidents(prs):
    s = slide(prs)
    hdr(s, "Real-World Hallucination Incidents",
        "What happens when AI outputs go unverified")

    incidents = [
        ("⚖  LEGAL",
         "Fabricated Case Citations",
         "Multiple attorneys sanctioned by federal courts for submitting AI-generated briefs "
         "containing case citations that do not exist. Some lost their bar licenses.",
         "Mata v. Avianca (SDNY 2023) — landmark case that triggered global AI policy debate.",
         RED),
        ("🏥  HEALTHCARE",
         "Dangerous Drug Dosing",
         "AI clinical tools recommended metformin starting at 1,000mg twice daily — double the "
         "FDA-approved 500mg starting dose. Also cited a grapefruit interaction with no clinical basis.",
         "Missed contraindication: Alcohol + metformin significantly increases lactic acidosis risk.",
         ORANGE),
        ("🏦  BANKING",
         "Fabricated Financial Terms",
         "Customer-facing chatbots quoted specific mortgage rates (5.25% APR), referenced promotions "
         "that did not exist in any product database, and guaranteed 24-hour pre-approval.",
         "All three elements are regulatory violations — banks cannot legally make such guarantees.",
         BLUE),
        ("🏛  GOVERNMENT",
         "Unauthorized Determinations",
         "AI systems issued benefits eligibility determinations not grounded in policy, "
         "and cited non-existent regulatory guidance when questioned by applicants.",
         "Citizens received incorrect benefit decisions based on fabricated policy interpretations.",
         PURPLE),
    ]

    for i, (ind, title, body, note, col) in enumerate(incidents):
        col_pos = i % 2
        row_pos = i // 2
        x = 0.35 + col_pos * 6.55
        y = 1.1 + row_pos * 2.95
        rect(s, x, y, 6.35, 2.75, CARD)
        rect(s, x, y, 0.06, 2.75, col)
        txt(s, ind, x + 0.2, y + 0.1, 5.8, 0.35, size=11, bold=True, color=col)
        txt(s, title, x + 0.2, y + 0.42, 5.8, 0.38, size=14, bold=True, color=WHITE)
        txt(s, body, x + 0.2, y + 0.82, 5.9, 1.0, size=10, color=LGRAY)
        rect(s, x + 0.2, y + 1.88, 5.9, 0.6, RGBColor(0x0A, 0x10, 0x20))
        txt(s, note, x + 0.3, y + 1.95, 5.7, 0.48, size=9, color=GRAY, italic=True)

# ============================================================================
# SLIDE 4 — ROOT CAUSE & MARKET GAP
# ============================================================================
def s04_root_cause(prs):
    s = slide(prs)
    hdr(s, "Root Cause & Market Gap", "Why this problem will not solve itself")

    # Conflict of interest
    txt(s, "THE STRUCTURAL CONFLICT OF INTEREST", 0.35, 1.1, 12.6, 0.35,
        size=13, bold=True, color=RED)
    rect(s, 0.35, 1.5, 12.6, 1.4, CARD)
    rect(s, 0.35, 1.5, 0.06, 1.4, RED)
    mtxt(s, [
        "AI providers (OpenAI, Google, Microsoft, Anthropic) have a fundamental commercial disincentive to",
        "solve hallucination at the output layer. Publishing failure rates and detection methodologies",
        "directly undermines confidence in their own products — reducing adoption and revenue.",
        "They built the problem. They cannot be trusted to solve it.",
    ], 0.6, 1.6, 12.1, 1.2, size=13, color=WHITE)

    # Why existing solutions fail
    txt(s, "WHY EXISTING SOLUTIONS FAIL", 0.35, 3.1, 12.6, 0.35, size=13, bold=True, color=ORANGE)
    failures = [
        ("Point Solutions",      "Address only narrow slices — citation checking OR bias detection, not both"),
        ("Governance Platforms", "Define policies but cannot actually detect hallucinations in real-time"),
        ("Provider Safety Filters","Designed for content moderation, not factual accuracy verification"),
        ("Custom In-House Builds","6–12 month build, $500K+, and ongoing maintenance burden"),
    ]
    for i, (sol, prob) in enumerate(failures):
        x = 0.35 + (i % 2) * 6.3
        y = 3.5 + (i // 2) * 1.0
        rect(s, x, y, 6.1, 0.85, CARD2)
        rect(s, x, y, 0.05, 0.85, ORANGE)
        txt(s, sol, x + 0.15, y + 0.08, 5.8, 0.32, size=12, bold=True, color=ORANGE)
        txt(s, prob, x + 0.15, y + 0.42, 5.8, 0.36, size=11, color=GRAY)

    # The gap
    rect(s, 0.35, 5.65, 12.6, 1.3, RGBColor(0x00, 0x22, 0x55))
    rect(s, 0.35, 5.65, 0.06, 1.3, BLUE)
    txt(s, "THE MARKET GAP", 0.55, 5.72, 4, 0.35, size=13, bold=True, color=BLUE)
    mtxt(s, [
        "No platform today combines real-time detection + automated prevention + governance",
        "infrastructure in a single, LLM-agnostic, enterprise-grade solution.",
        "That gap is exactly what TrustLayer AI fills.",
    ], 0.55, 6.1, 12.1, 0.8, size=12, color=WHITE)

# ============================================================================
# SLIDE 5 — SOLUTION OVERVIEW
# ============================================================================
def s05_solution(prs):
    s = slide(prs)
    hdr(s, "TrustLayer AI — Solution Overview",
        "The enterprise trust layer between your business and any AI system")

    # Pipeline
    boxes = [
        ("User\nQuery",      CARD,  WHITE, "Any channel:\nchat, API, app"),
        ("LLM\nProvider",    CARD,  WHITE, "OpenAI, Anthropic,\nGoogle, Azure…"),
        ("TrustLayer\nAI",   BLUE,  WHITE, "Intercept · Analyze\nScore · Act"),
        ("Verified\nResponse",TEAL, WHITE, "Delivered or\nblocked/flagged"),
    ]
    for i, (label, bg, fg, sub) in enumerate(boxes):
        x = 0.35 + i * 3.15
        rect(s, x, 1.1, 2.85, 1.1, bg)
        txt(s, label, x, 1.1, 2.85, 0.75, size=15, bold=True, color=fg, align=PP_ALIGN.CENTER)
        txt(s, sub, x + 0.1, 1.8, 2.65, 0.38, size=9, color=GRAY, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, "→", x + 2.85, 1.38, 0.3, 0.45, size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Three outcomes
    outcomes = [
        ("PASS",  TEAL,   "Confidence ≥ 75%\nRisk Score < 30",
         "Response is verified and reliable.\nDelivered directly to the user.",
         ["Accuracy confirmed", "Latency < 100ms", "Full audit logged"]),
        ("FLAG",  ORANGE, "Confidence 50–74%\nRisk Score 30–59",
         "Response is uncertain.\nRouted to a human reviewer before delivery.",
         ["Review queue created", "Reviewer notified", "SLA tracked"]),
        ("BLOCK", RED,    "Confidence < 50%\nRisk Score ≥ 60",
         "Response is dangerous or unverifiable.\nBlocked; fallback action triggered automatically.",
         ["User protected", "Incident logged", "Compliance flagged"]),
    ]
    for i, (action, col, criteria, desc, bullets) in enumerate(outcomes):
        x = 0.35 + i * 4.3
        rect(s, x, 2.45, 4.1, 4.7, CARD)
        rect(s, x, 2.45, 4.1, 0.5, col)
        txt(s, action, x, 2.45, 4.1, 0.5, size=20, bold=True, align=PP_ALIGN.CENTER)
        txt(s, criteria, x + 0.15, 3.05, 3.8, 0.5, size=11, bold=True, color=col)
        txt(s, desc, x + 0.15, 3.62, 3.8, 0.9, size=11, color=GRAY)
        for j, b in enumerate(bullets):
            rect(s, x + 0.15, 4.6 + j * 0.58, 3.8, 0.48, CARD2)
            rect(s, x + 0.15, 4.6 + j * 0.58, 0.04, 0.48, col)
            txt(s, b, x + 0.3, 4.68 + j * 0.58, 3.6, 0.32, size=11, color=WHITE)

# ============================================================================
# SLIDE 6 — SYSTEM ARCHITECTURE DIAGRAM
# ============================================================================
def s06_architecture(prs):
    s = slide(prs)
    hdr(s, "System Architecture",
        "Request interception → parallel detection → scored decision → audit")

    # Layer labels
    layers = [
        (1.05,  "ENTERPRISE LAYER",    BLUE),
        (2.25,  "TRUSTLAYER GATEWAY",  TEAL),
        (3.45,  "DETECTION ENGINE",    ORANGE),
        (5.25,  "DATA SOURCES",        PURPLE),
        (6.35,  "ACTIONS",             GREEN),
    ]
    for y, label, col in layers:
        rect(s, 0.1, y, 1.6, 0.9, CARD2)
        txt(s, label, 0.12, y + 0.1, 1.55, 0.7, size=8, bold=True, color=col,
            align=PP_ALIGN.CENTER)

    # Enterprise apps box
    rect(s, 1.85, 1.05, 3.5, 0.9, CARD)
    txt(s, "Enterprise Applications", 1.85, 1.05, 3.5, 0.45, size=12, bold=True,
        align=PP_ALIGN.CENTER)
    txt(s, "Chatbots  ·  Copilots  ·  Internal Tools  ·  APIs",
        1.85, 1.45, 3.5, 0.4, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # LLM Providers box
    rect(s, 6.2, 1.05, 3.5, 0.9, CARD)
    txt(s, "LLM Providers", 6.2, 1.05, 3.5, 0.45, size=12, bold=True, align=PP_ALIGN.CENTER)
    txt(s, "OpenAI · Anthropic · Google · Azure · AWS · Mistral",
        6.2, 1.45, 3.5, 0.4, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Arrows down
    for x in [3.35, 7.85]:
        txt(s, "↓", x, 1.95, 0.5, 0.3, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    # Gateway
    rect(s, 1.85, 2.25, 7.85, 0.95, DARK_BLUE)
    rect(s, 1.85, 2.25, 7.85, 0.08, TEAL)
    txt(s, "TrustLayer API Gateway  —  Zero-code interception; LLM-agnostic routing",
        1.85, 2.35, 7.85, 0.75, size=13, bold=True, align=PP_ALIGN.CENTER)

    txt(s, "↓", 5.6, 3.2, 0.5, 0.3, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # Detection Engine
    rect(s, 1.85, 3.45, 7.85, 1.75, CARD)
    rect(s, 1.85, 3.45, 7.85, 0.1, ORANGE)
    txt(s, "Detection Engine  (8 algorithms in parallel)", 1.85, 3.55, 7.85, 0.4,
        size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    algos = ["Semantic\nEntropy", "Self-\nConsistency", "Citation\nVerify",
             "Enterprise\nGrounding", "Claim\nExtraction", "Pattern\nRecognition",
             "Temporal\nCheck", "Numerical\nValidation"]
    for i, a in enumerate(algos):
        x = 1.95 + i * 0.96
        rect(s, x, 4.0, 0.88, 1.0, CARD2)
        txt(s, a, x, 4.05, 0.88, 0.9, size=8, color=LGRAY, align=PP_ALIGN.CENTER)

    # Data Sources (right column)
    rect(s, 10.2, 3.45, 3.0, 1.75, CARD)
    rect(s, 10.2, 3.45, 3.0, 0.1, PURPLE)
    txt(s, "Data Sources", 10.2, 3.55, 3.0, 0.3, size=11, bold=True, color=PURPLE,
        align=PP_ALIGN.CENTER)
    sources = ["CRM / Product Catalog", "Policy Databases", "Legal Repositories",
               "FDA / Clinical DBs", "Market Data Feeds"]
    for i, src in enumerate(sources):
        txt(s, "· " + src, 10.3, 3.92 + i * 0.24, 2.8, 0.22, size=9, color=GRAY)

    # Connector arrow
    txt(s, "←", 9.7, 4.3, 0.5, 0.3, size=16, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    # Scoring
    txt(s, "↓", 5.6, 5.2, 0.5, 0.3, size=18, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Score + decision box
    rect(s, 1.85, 5.5, 7.85, 0.72, RGBColor(0x0A, 0x14, 0x28))
    txt(s, "Scoring Model  →  Confidence Score (0–100)  +  Risk Score (0–100)  →  Pass / Flag / Block",
        1.85, 5.6, 7.85, 0.5, size=11, bold=True, align=PP_ALIGN.CENTER, color=LGRAY)

    # Actions
    txt(s, "↓", 5.6, 6.22, 0.5, 0.3, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    for i, (action, col, w) in enumerate([("PASS → User", GREEN, 2.5),
                                           ("FLAG → Human Review", ORANGE, 3.0),
                                           ("BLOCK → Fallback", RED, 2.8)]):
        x = 1.85 + i * 3.2 + (0.2 if i else 0)
        rect(s, x, 6.52, w, 0.6, col)
        txt(s, action, x, 6.52, w, 0.6, size=11, bold=True, align=PP_ALIGN.CENTER)

    # Audit note
    rect(s, 10.2, 5.5, 3.0, 1.62, CARD)
    txt(s, "Audit & Compliance", 10.2, 5.55, 3.0, 0.3, size=10, bold=True,
        color=TEAL, align=PP_ALIGN.CENTER)
    audit_items = ["Immutable audit log", "Compliance tagging", "Dashboard reporting",
                   "Auto-documentation", "Incident alerts"]
    for i, item in enumerate(audit_items):
        txt(s, "· " + item, 10.3, 5.92 + i * 0.23, 2.8, 0.22, size=9, color=GRAY)

# ============================================================================
# SLIDE 7 — DETECTION ENGINE OVERVIEW
# ============================================================================
def s07_detection_overview(prs):
    s = slide(prs)
    hdr(s, "Detection Engine — 8-Technique Architecture",
        "Parallel algorithms targeting distinct LLM failure modes")

    # Table header
    cols = ["Technique", "Accuracy", "Latency", "Failure Mode Targeted"]
    widths = [3.8, 1.1, 1.0, 6.0]
    xs = [0.35, 4.2, 5.35, 6.4]
    for j, (h, x, w) in enumerate(zip(cols, xs, widths)):
        rect(s, x, 1.1, w, 0.4, BLUE)
        txt(s, h, x + 0.05, 1.1, w - 0.1, 0.4, size=11, bold=True, align=PP_ALIGN.CENTER)

    techniques = [
        ("Semantic Entropy Analysis",        "87%", "45ms",  "Internal model uncertainty; knowledge boundary violations"),
        ("Self-Consistency Checking",        "82%", "120ms", "Cross-query contradictions and logical inconsistencies"),
        ("Source Citation Verification",     "95%", "200ms", "Fabricated URLs, case citations, document references"),
        ("Enterprise Knowledge Grounding",   "98%", "150ms", "Claims contradicting authoritative enterprise data"),
        ("Claim Extraction & Classification","91%", "35ms",  "High-risk assertion types: financial, medical, legal"),
        ("Hallucination Pattern Recognition","79%", "25ms",  "Known LLM failure signatures from curated datasets"),
        ("Temporal Consistency Checking",    "93%", "40ms",  "Outdated or temporally incoherent information"),
        ("Numerical Claim Validation",       "96%", "80ms",  "Fabricated rates, dosages, thresholds, percentages"),
    ]
    for i, (name, acc, lat, desc) in enumerate(techniques):
        y = 1.55 + i * 0.62
        row_bg = CARD if i % 2 == 0 else CARD2
        for j, (x, w) in enumerate(zip(xs, widths)):
            rect(s, x, y, w, 0.57, row_bg)
        txt(s, name,  xs[0] + 0.1, y + 0.1, widths[0] - 0.15, 0.4, size=11, color=WHITE)
        txt(s, acc,   xs[1] + 0.05, y + 0.1, widths[1] - 0.1, 0.4, size=12, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        txt(s, lat,   xs[2] + 0.05, y + 0.1, widths[2] - 0.1, 0.4, size=11, color=ORANGE, align=PP_ALIGN.CENTER)
        txt(s, desc,  xs[3] + 0.08, y + 0.1, widths[3] - 0.12, 0.4, size=10, color=GRAY)

    # Summary bar
    rect(s, 0.35, 6.55, 12.6, 0.7, BLUE)
    summary = [("92%", "Overall catch rate"), ("<3%", "False positive rate"),
               ("<100ms", "Pipeline latency"), ("99.9%", "Platform SLA")]
    for i, (val, lbl) in enumerate(summary):
        x = 0.7 + i * 3.15
        txt(s, val, x, 6.62, 2.0, 0.38, size=20, bold=True, align=PP_ALIGN.CENTER)
        txt(s, lbl, x, 6.97, 2.0, 0.22, size=10, color=LBLUE, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 8 — DETECTION DEEP-DIVE: GROUNDING + ENTROPY
# ============================================================================
def s08_detection_deepdive(prs):
    s = slide(prs)
    hdr(s, "Detection Deep-Dive — Two Highest-Accuracy Techniques",
        "Enterprise Grounding (98%) and Source Citation Verification (95%)")

    # Grounding
    rect(s, 0.35, 1.1, 6.2, 5.8, CARD)
    rect(s, 0.35, 1.1, 6.2, 0.55, TEAL)
    txt(s, "Enterprise Knowledge Grounding  —  98% accuracy", 0.35, 1.1, 6.2, 0.55,
        size=13, bold=True, align=PP_ALIGN.CENTER)

    mtxt(s, [
        "What it does:",
        "Extracts all factual claims from the AI response — prices, product specs,",
        "policies, personnel details, inventory, rates — and validates each claim",
        "against your live authoritative data sources in real-time.",
    ], 0.5, 1.75, 6.0, 1.1, size=11, color=WHITE)

    mtxt(s, [
        "How it works:",
        "1. NLP claim extraction identifies all verifiable assertions",
        "2. Each claim is mapped to the relevant enterprise data source",
        "3. Real-time API call to CRM / ERP / policy DB / product catalog",
        "4. Response score drops for every unverified or contradicted claim",
    ], 0.5, 2.95, 6.0, 1.4, size=11, color=LGRAY)

    mtxt(s, [
        "Why 98% accuracy:",
        "Unlike probabilistic approaches, grounding compares against ground truth.",
        "If the AI says '5.25% APR' and the product DB says '6.1% APR', that is a",
        "definitive failure — no ambiguity, no false positives.",
    ], 0.5, 4.45, 6.0, 1.1, size=11, color=WHITE)

    tag(s, 0.5, 5.65, "150ms latency", TEAL, w=1.4, h=0.32)
    tag(s, 2.1, 5.65, "Requires enterprise connector", DARK_BLUE, w=2.5, h=0.32)
    tag(s, 4.75, 5.65, "Highest value technique", GREEN, w=1.7, h=0.32)

    # Citation verification
    rect(s, 6.95, 1.1, 6.05, 5.8, CARD)
    rect(s, 6.95, 1.1, 6.05, 0.55, BLUE)
    txt(s, "Source Citation Verification  —  95% accuracy", 6.95, 1.1, 6.05, 0.55,
        size=13, bold=True, align=PP_ALIGN.CENTER)

    mtxt(s, [
        "What it does:",
        "Extracts every citation, URL, case reference, statute, document name,",
        "and publication the AI response mentions — then attempts real-time",
        "verification of each one against indexed databases.",
    ], 7.1, 1.75, 5.85, 1.1, size=11, color=WHITE)

    mtxt(s, [
        "Verification sources:",
        "· Legal:     Westlaw, LexisNexis, CourtListener, Federal Register",
        "· Medical:   PubMed, FDA databases, clinical guidelines",
        "· Web:       Real-time URL resolution and content matching",
        "· Academic:  CrossRef, DOI lookup, JSTOR",
    ], 7.1, 2.95, 5.85, 1.4, size=11, color=LGRAY)

    mtxt(s, [
        "Real example — Legal hallucination:",
        'AI cited "Morrison v. Tesla Motors, 9th Cir. 2023" as precedent.',
        "TrustLayer searched Westlaw, LexisNexis, CourtListener.",
        "Case does not exist. Confidence dropped to 4%. Response BLOCKED.",
        "Attorney saved from potential sanctions and bar discipline.",
    ], 7.1, 4.45, 5.85, 1.35, size=11, color=WHITE)

    tag(s, 7.1, 5.65, "200ms latency", BLUE, w=1.4, h=0.32)
    tag(s, 8.65, 5.65, "Real-time API lookups", DARK_BLUE, w=2.0, h=0.32)
    tag(s, 10.8, 5.65, "Critical for Legal & Medical", RED, w=2.1, h=0.32)

# ============================================================================
# SLIDE 9 — SCORING MODEL
# ============================================================================
def s09_scoring(prs):
    s = slide(prs)
    hdr(s, "Scoring Model & Decision Logic",
        "How 8 technique scores combine into a single Pass / Flag / Block decision")

    # Formula
    rect(s, 0.35, 1.1, 12.6, 1.6, CARD)
    rect(s, 0.35, 1.1, 0.06, 1.6, BLUE)
    txt(s, "Confidence Score  =  weighted_average(8 technique scores)  ×  grounding_multiplier",
        0.55, 1.18, 12.2, 0.45, size=13, bold=True, color=BLUE)
    txt(s, "Risk Score  =  severity_weight  ×  (1 − confidence)  ×  domain_risk_factor",
        0.55, 1.62, 12.2, 0.45, size=13, bold=True, color=ORANGE)
    txt(s, "domain_risk_factor is higher for Healthcare and Legal (life / liberty at stake)",
        0.55, 2.06, 12.2, 0.4, size=11, color=GRAY, italic=True)

    # Decision matrix
    txt(s, "DECISION MATRIX", 0.35, 2.85, 12.6, 0.35, size=12, bold=True, color=WHITE)
    rows = [
        ("PASS",  "≥ 75",  "< 30",  "Response delivered immediately to user",       TEAL),
        ("FLAG",  "50–74", "30–59", "Queued for human review; SLA tracked",          ORANGE),
        ("BLOCK", "< 50",  "≥ 60",  "Response blocked; fallback triggered; incident logged", RED),
    ]
    hdrs = ["Action", "Confidence", "Risk Score", "Outcome"]
    hdr_widths = [1.5, 1.8, 1.8, 7.0]
    hdr_x = [0.35, 1.9, 3.75, 5.6]
    for j, (h, x, w) in enumerate(zip(hdrs, hdr_x, hdr_widths)):
        rect(s, x, 3.25, w, 0.4, DARK_BLUE)
        txt(s, h, x + 0.05, 3.25, w - 0.1, 0.4, size=11, bold=True, align=PP_ALIGN.CENTER)
    for i, (action, conf, risk, outcome, col) in enumerate(rows):
        y = 3.7 + i * 0.75
        for j, (x, w) in enumerate(zip(hdr_x, hdr_widths)):
            rect(s, x, y, w, 0.68, CARD2 if i % 2 == 0 else CARD)
        rect(s, hdr_x[0], y, hdr_widths[0], 0.68, col)
        txt(s, action, hdr_x[0], y + 0.1, hdr_widths[0], 0.45, size=14, bold=True, align=PP_ALIGN.CENTER)
        txt(s, conf,   hdr_x[1] + 0.1, y + 0.15, hdr_widths[1] - 0.2, 0.4, size=12, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, risk,   hdr_x[2] + 0.1, y + 0.15, hdr_widths[2] - 0.2, 0.4, size=12, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, outcome, hdr_x[3] + 0.1, y + 0.15, hdr_widths[3] - 0.2, 0.4, size=11, color=LGRAY)

    # Latency guarantee
    rect(s, 0.35, 6.05, 12.6, 1.1, RGBColor(0x00, 0x22, 0x55))
    rect(s, 0.35, 6.05, 0.06, 1.1, TEAL)
    txt(s, "Latency Guarantee", 0.55, 6.1, 4, 0.35, size=13, bold=True, color=TEAL)
    mtxt(s, [
        "Total TrustLayer overhead: < 100ms (8 techniques run in parallel, not sequentially)",
        "API Gateway routing: <2ms  |  Detection engine: <80ms  |  Scoring & decision: <5ms  |  Delivery: <5ms",
    ], 0.55, 6.45, 12.1, 0.62, size=11, color=WHITE)

# ============================================================================
# SLIDE 10 — INDUSTRY MODULES OVERVIEW
# ============================================================================
def s10_industries(prs):
    s = slide(prs)
    hdr(s, "Industry Modules — 8 Specialized Verticals",
        "Generic detection is not enough. Each industry has distinct risk profiles and data sources.")

    modules = [
        ("BFSI",          BLUE,   "Banking & Financial Services",
         "Fabricated rates, unauthorized eligibility, regulatory disclosures"),
        ("Healthcare",    RED,    "Clinical & Administrative",
         "Drug dosing errors, dangerous contraindications, unauthorized approvals"),
        ("Legal",         ORANGE, "Litigation & Contract",
         "Fabricated citations, non-existent statutes, unauthorized determinations"),
        ("Enterprise",    TEAL,   "HR & IT Operations",
         "Fabricated HR policies, incorrect IT guidance, process hallucinations"),
        ("Manufacturing", PURPLE, "Quality & Compliance",
         "Fabricated QC records, ISO certification claims, spec violations"),
        ("Retail",        GREEN,  "Product & Pricing",
         "Unverified pricing, inventory errors, promotional term fabrication"),
        ("Telecom",       RGBColor(0x27,0x6E,0xBD), "Billing & Account",
         "Unauthorized account actions, billing errors, plan misrepresentation"),
        ("Government",    RGBColor(0xC0,0x39,0x2B), "Benefits & Policy",
         "Unauthorized eligibility, policy misinterpretation, regulatory fabrication"),
    ]

    for i, (code, col, title, desc) in enumerate(modules):
        col_n = i % 4
        row_n = i // 4
        x = 0.35 + col_n * 3.25
        y = 1.1 + row_n * 2.1
        rect(s, x, y, 3.05, 1.9, CARD)
        rect(s, x, y, 3.05, 0.48, col)
        txt(s, code, x, y, 2.3, 0.48, size=15, bold=True, align=PP_ALIGN.CENTER)
        txt(s, title, x + 0.1, y + 0.55, 2.85, 0.38, size=11, bold=True, color=WHITE)
        txt(s, desc, x + 0.1, y + 0.95, 2.85, 0.88, size=10, color=GRAY)
        tag(s, x + 2.1, y + 0.08, "MODULE", col, w=0.88, h=0.3)

    # Bottom: what each module includes
    rect(s, 0.35, 5.35, 12.6, 1.85, CARD2)
    rect(s, 0.35, 5.35, 0.06, 1.85, BLUE)
    txt(s, "Every module includes:", 0.55, 5.42, 4, 0.35, size=12, bold=True, color=BLUE)
    items = [
        ("Industry-specific risk models", "Pre-trained on domain hallucination datasets"),
        ("Pre-built compliance rule sets", "Maps to relevant regulations per vertical"),
        ("Enterprise system integrations", "Pre-built connectors to authoritative data sources"),
        ("Domain-specific scoring weights","Higher penalties for life-safety critical errors"),
    ]
    for i, (title, desc) in enumerate(items):
        x = 0.55 + (i % 2) * 6.3
        y = 5.8 + (i // 2) * 0.62
        txt(s, "→  " + title + ":", x, y, 6.0, 0.28, size=11, bold=True, color=WHITE)
        txt(s, "     " + desc, x, y + 0.28, 6.0, 0.28, size=10, color=GRAY)

# ============================================================================
# SLIDE 11 — BFSI DEEP-DIVE (BANKING DEMO WALKTHROUGH)
# ============================================================================
def s11_bfsi(prs):
    s = slide(prs)
    hdr(s, "BFSI Module — Banking Demo Walkthrough",
        "Step-by-step: how TrustLayer catches a mortgage hallucination")

    # User query
    txt(s, "STEP 1 — Customer Query", 0.35, 1.1, 12.6, 0.32, size=11, bold=True, color=GRAY)
    rect(s, 0.35, 1.45, 12.6, 0.72, CARD2)
    txt(s, '"What is the current interest rate for a 30-year fixed mortgage, and am I pre-approved?"',
        0.55, 1.55, 12.2, 0.55, size=13, color=WHITE, italic=True)

    # AI response
    txt(s, "STEP 2 — Raw LLM Response (before TrustLayer)", 0.35, 2.28, 12.6, 0.32,
        size=11, bold=True, color=RED)
    rect(s, 0.35, 2.62, 12.6, 1.0, RGBColor(0x2A, 0x10, 0x10))
    rect(s, 0.35, 2.62, 0.06, 1.0, RED)
    mtxt(s, [
        '"Great news! Our current rate for a 30-year fixed mortgage is 5.25% APR.',
        'As an existing customer, you qualify for our Spring Promotion with 0.5% rate reduction.',
        'Based on your profile, you are guaranteed pre-approval within 24 hours."',
    ], 0.55, 2.72, 12.2, 0.82, size=12, color=WHITE, italic=True)

    # Issues detected
    txt(s, "STEP 3 — TrustLayer Detection Results", 0.35, 3.73, 12.6, 0.32,
        size=11, bold=True, color=ORANGE)
    issues = [
        ("Fabricated Rate",         "5.25% APR does not exist in live pricing feed — actual rate is 6.85%",      RED),
        ("Non-existent Promotion",  '"Spring Promotion" not found in marketing database (0 records)',             RED),
        ("Illegal Guarantee",       "Banks cannot legally guarantee approval — violates TILA / Reg Z",           RED),
        ("Fiduciary Violation",     "Providing specific rate advice without disclosures required by law",        ORANGE),
        ("Unsupported Claim",       "Profile assessment made without access to applicant credit data",           ORANGE),
    ]
    for i, (issue, detail, col) in enumerate(issues):
        y = 4.1 + i * 0.42
        row_bg = CARD if i % 2 == 0 else CARD2
        rect(s, 0.35, y, 12.6, 0.38, row_bg)
        rect(s, 0.35, y, 0.04, 0.38, col)
        txt(s, issue + ":", 0.5, y + 0.05, 2.5, 0.28, size=10, bold=True, color=col)
        txt(s, detail, 3.1, y + 0.05, 9.7, 0.28, size=10, color=LGRAY)

    # Final scores
    rect(s, 0.35, 6.25, 12.6, 0.95, BLUE)
    txt(s, "STEP 4 — Decision", 0.55, 6.3, 2.0, 0.3, size=11, bold=True, color=LBLUE)
    for val, lbl, x in [
        ("18%",      "Confidence Score", 2.8),
        ("94%",      "Risk Score",       5.5),
        ("BLOCKED",  "Action",           8.1),
        ("Specialist","Routed to",       10.4),
    ]:
        txt(s, val, x, 6.3, 2.0, 0.45, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, lbl, x, 6.72, 2.0, 0.3, size=10, color=LBLUE, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 12 — HEALTHCARE DEEP-DIVE
# ============================================================================
def s12_healthcare(prs):
    s = slide(prs)
    hdr(s, "Healthcare Module — Drug Dosing Hallucination",
        "A single incorrect AI response here can be life-threatening")

    # Clinical question
    rect(s, 0.35, 1.1, 12.6, 0.68, CARD2)
    txt(s, 'Clinical Query: "What is the correct starting dose of metformin for a newly-diagnosed Type 2 diabetic?"',
        0.55, 1.2, 12.2, 0.52, size=13, color=WHITE, italic=True)

    # AI hallucinated response
    txt(s, "AI RESPONSE — 3 DANGEROUS ERRORS DETECTED", 0.35, 1.92, 12.6, 0.32,
        size=11, bold=True, color=RED)
    rect(s, 0.35, 2.27, 12.6, 1.2, RGBColor(0x2A, 0x10, 0x10))
    mtxt(s, [
        '"Start at 1,000mg twice daily. Take with grapefruit juice to enhance absorption.',
        ' There are no concerns with alcohol consumption while taking metformin."',
    ], 0.55, 2.37, 12.2, 0.6, size=13, color=WHITE, italic=True)

    # Three error cards
    errors = [
        ("ERROR 1\nDosing",
         "AI recommends: 1,000mg twice daily\nFDA approved starting dose: 500mg once daily\n\n"
         "Recommending double the starting dose risks GI toxicity and patient harm.",
         RED),
        ("ERROR 2\nFabricated Interaction",
         "Grapefruit juice interaction: NO clinical evidence exists.\n"
         "AI completely fabricated this pharmacokinetic claim.\n\n"
         "Source: FDA drug label — no grapefruit interaction listed.",
         ORANGE),
        ("ERROR 3\nMissed Contraindication",
         "Alcohol + metformin significantly increases risk of lactic acidosis.\n"
         "Lactic acidosis is potentially fatal.\n\n"
         "This contraindication is in every clinical guideline.",
         RED),
    ]
    for i, (title, detail, col) in enumerate(errors):
        x = 0.35 + i * 4.3
        rect(s, x, 3.6, 4.1, 2.9, CARD)
        rect(s, x, 3.6, 4.1, 0.55, col)
        txt(s, title, x, 3.6, 4.1, 0.55, size=13, bold=True, align=PP_ALIGN.CENTER)
        txt(s, detail, x + 0.12, 4.22, 3.86, 2.2, size=10.5, color=LGRAY)

    # Result
    rect(s, 0.35, 6.62, 12.6, 0.68, RGBColor(0x2A, 0x10, 0x10))
    rect(s, 0.35, 6.62, 12.6, 0.06, RED)
    for val, lbl, x in [("6%", "Confidence Score", 0.6), ("98%", "Risk Score", 3.4),
                         ("BLOCKED", "Action", 6.0), ("Clinical Pharmacist", "Escalated to", 8.8)]:
        txt(s, val, x, 6.7, 2.3, 0.38, size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
        txt(s, lbl, x, 7.05, 2.3, 0.2, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 13 — LEGAL DEEP-DIVE
# ============================================================================
def s13_legal(prs):
    s = slide(prs)
    hdr(s, "Legal Module — Fabricated Case Citations",
        "The most publicly documented AI hallucination failure — career-ending if undetected")

    # Query
    rect(s, 0.35, 1.1, 12.6, 0.65, CARD2)
    txt(s, '"Research case law on autonomous vehicle manufacturer liability for software failures."',
        0.55, 1.2, 12.2, 0.48, size=13, color=WHITE, italic=True)

    # Fabricated citations
    txt(s, "AI-GENERATED CITATIONS — ALL FABRICATED", 0.35, 1.88, 12.6, 0.32,
        size=11, bold=True, color=RED)

    citations = [
        ("Morrison v. Tesla Motors, Inc.",     "9th Circuit Court of Appeals, 2023",  "Case No. 22-16598"),
        ("Waymo LLC v. California DMV",        "U.S. Supreme Court, 2022",            "Case No. 21-1024"),
        ("National Transportation Safety AI Division v. Ford", "D.C. Circuit, 2023",  "Case No. 23-1147"),
        ("AI Liability Act of 2023",           "United States Congress",              "Pub. L. 118-22"),
    ]
    for i, (case, court, ref) in enumerate(citations):
        y = 2.3 + i * 0.68
        rect(s, 0.35, y, 12.6, 0.6, CARD if i % 2 == 0 else CARD2)
        rect(s, 0.35, y, 0.06, 0.6, RED)
        txt(s, case, 0.55, y + 0.05, 7.0, 0.3, size=12, bold=True, color=WHITE)
        txt(s, court + "  ·  " + ref, 0.55, y + 0.33, 8.0, 0.22, size=10, color=GRAY)
        tag(s, 10.0, y + 0.14, "NOT FOUND", RED, w=1.8, h=0.3)
        tag(s, 8.1, y + 0.14, "FABRICATED", ORANGE, w=1.8, h=0.3)

    # Verification process
    txt(s, "TRUSTLAYER VERIFICATION PROCESS", 0.35, 5.08, 12.6, 0.32,
        size=11, bold=True, color=BLUE)
    steps = [
        ("Extract", "NLP extracts all citation entities from AI response"),
        ("Search",  "Parallel queries: Westlaw, LexisNexis, CourtListener, Federal Register"),
        ("Score",   "Each unresolved citation reduces confidence score significantly"),
        ("Block",   "4% confidence → BLOCKED; attorney saved from potential sanctions"),
    ]
    for i, (step, desc) in enumerate(steps):
        x = 0.35 + i * 3.22
        rect(s, x, 5.45, 3.05, 1.15, CARD)
        rect(s, x, 5.45, 3.05, 0.38, BLUE)
        txt(s, step, x, 5.45, 3.05, 0.38, size=14, bold=True, align=PP_ALIGN.CENTER)
        txt(s, desc, x + 0.12, 5.9, 2.82, 0.62, size=10, color=GRAY)

    rect(s, 0.35, 6.72, 12.6, 0.58, RED)
    txt(s, "Confidence: 4%  ·  Action: BLOCKED  ·  All 4 citations fabricated  ·  "
        "Saved from court sanctions, malpractice, and possible bar discipline",
        0.5, 6.82, 12.2, 0.4, size=12, bold=True, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 14 — COMPLIANCE & GOVERNANCE
# ============================================================================
def s14_compliance(prs):
    s = slide(prs)
    hdr(s, "Compliance & Governance",
        "One platform — four major regulatory frameworks covered out of the box")

    frameworks = [
        ("EU AI Act",     "94%", "Aug 2026",
         ["Article 9: Risk management system", "Article 11: Technical documentation",
          "Article 14: Human oversight mechanisms", "Article 72: Accuracy monitoring & logging"],
         "Penalties: up to €35M or 7% global revenue", RED),
        ("NIST AI RMF",   "89%", "US standard",
         ["GOVERN: Policies, accountability, culture", "MAP: Categorize AI risks",
          "MEASURE: Analyze and assess risks", "MANAGE: Prioritize and treat risks"],
         "Referenced in US federal AI procurement requirements", BLUE),
        ("ISO/IEC 42001", "82%", "International",
         ["AI management system documentation", "Risk assessment framework",
          "AI objectives and performance evaluation", "Continual improvement processes"],
         "Emerging international standard for AI governance", TEAL),
        ("SOC 2 Type II", "97%", "Ongoing",
         ["Availability: 99.9% uptime SLA", "Processing integrity: accurate outputs",
          "Confidentiality: data protection controls", "Security: access controls and monitoring"],
         "Required by many enterprise vendor evaluation processes", GREEN),
    ]

    for i, (fw, pct, period, reqs, note, col) in enumerate(frameworks):
        col_n = i % 2
        row_n = i // 2
        x = 0.35 + col_n * 6.55
        y = 1.1 + row_n * 2.95
        rect(s, x, y, 6.25, 2.72, CARD)
        rect(s, x, y, 6.25, 0.52, col)
        txt(s, fw, x + 0.12, y + 0.06, 3.0, 0.42, size=15, bold=True)
        txt(s, pct + " coverage", x + 3.5, y + 0.1, 2.5, 0.32, size=18, bold=True,
            color=col, align=PP_ALIGN.RIGHT)
        txt(s, period, x + 4.8, y + 0.38, 1.3, 0.2, size=9, color=LBLUE, align=PP_ALIGN.RIGHT)

        # Progress bar
        bar_w = 5.8 * int(pct.replace("%", "")) / 100
        rect(s, x + 0.15, y + 0.65, 5.8, 0.2, CARD2)
        rect(s, x + 0.15, y + 0.65, bar_w, 0.2, col)

        for j, req in enumerate(reqs):
            txt(s, "·  " + req, x + 0.15, y + 0.95 + j * 0.3, 5.9, 0.28, size=10, color=LGRAY)
        rect(s, x + 0.15, y + 2.3, 5.9, 0.28, RGBColor(0x0A, 0x14, 0x28))
        txt(s, note, x + 0.25, y + 2.34, 5.7, 0.2, size=9, color=GRAY, italic=True)

# ============================================================================
# SLIDE 15 — EU AI ACT DEEP-DIVE
# ============================================================================
def s15_eu_ai_act(prs):
    s = slide(prs)
    hdr(s, "EU AI Act — August 2026 Enforcement",
        "The most significant AI regulation globally — TrustLayer provides 94% coverage out of the box")

    # Timeline
    txt(s, "ENFORCEMENT TIMELINE", 0.35, 1.1, 12.6, 0.3, size=11, bold=True, color=GRAY)
    milestones = [
        ("Feb 2025", "Prohibited AI provisions in effect"),
        ("Aug 2025", "GPAI model rules apply"),
        ("Aug 2026", "HIGH-RISK AI FULL ENFORCEMENT"),
        ("Aug 2027", "Embedded AI systems deadline"),
    ]
    for i, (date, label) in enumerate(milestones):
        x = 0.5 + i * 3.15
        col = RED if "FULL" in label else BLUE
        rect(s, x, 1.45, 2.9, 0.6, CARD)
        rect(s, x, 1.45, 2.9, 0.08, col)
        txt(s, date, x + 0.1, 1.55, 2.7, 0.28, size=12, bold=True, color=col)
        txt(s, label, x + 0.1, 1.82, 2.7, 0.2, size=9, color=GRAY)

    txt(s, "PENALTIES FOR NON-COMPLIANCE", 0.35, 2.2, 12.6, 0.3, size=11, bold=True, color=RED)
    rect(s, 0.35, 2.55, 12.6, 0.62, RGBColor(0x2A, 0x10, 0x10))
    rect(s, 0.35, 2.55, 12.6, 0.07, RED)
    txt(s, "Up to €35,000,000  or  7% of total worldwide annual turnover (whichever is higher)",
        0.45, 2.68, 12.3, 0.38, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # 5 requirements
    txt(s, "5 CORE REQUIREMENTS — TRUSTLAYER COVERAGE", 0.35, 3.3, 12.6, 0.3,
        size=11, bold=True, color=TEAL)
    reqs = [
        ("Risk Classification",
         "Article 6 — Classify AI systems by risk tier (Unacceptable, High, Limited, Minimal)",
         "TrustLayer auto-classifies each deployment. Dashboard shows risk tier per AI system."),
        ("Technical Documentation",
         "Article 11 — Maintain technical documentation before market placement",
         "Auto-generated for every deployment. No manual assembly required."),
        ("Human Oversight",
         "Article 14 — Implement effective human oversight measures",
         "Every FLAG and BLOCK creates an auditable human review record."),
        ("Accuracy & Robustness",
         "Article 15 — Achieve appropriate levels of accuracy, robustness, cybersecurity",
         "Real-time accuracy monitoring dashboard. 92% catch rate with <3% false positives."),
        ("Post-Market Monitoring",
         "Article 72 — Establish post-market monitoring system",
         "Continuous hallucination trend monitoring. Automated incident reporting to authorities."),
    ]
    for i, (title, article, impl) in enumerate(reqs):
        y = 3.68 + i * 0.67
        row_bg = CARD if i % 2 == 0 else CARD2
        rect(s, 0.35, y, 12.6, 0.6, row_bg)
        rect(s, 0.35, y, 0.05, 0.6, TEAL)
        txt(s, title + ":", 0.52, y + 0.05, 2.5, 0.28, size=11, bold=True, color=TEAL)
        txt(s, article, 0.52, y + 0.32, 5.5, 0.24, size=9, color=GRAY)
        txt(s, impl, 6.2, y + 0.12, 6.6, 0.38, size=10, color=WHITE)

# ============================================================================
# SLIDE 16 — INTEGRATION ECOSYSTEM
# ============================================================================
def s16_integrations(prs):
    s = slide(prs)
    hdr(s, "Integration Ecosystem",
        "LLM-agnostic  ·  Zero vendor lock-in  ·  8 providers  ·  50+ enterprise connectors")

    # LLM providers
    txt(s, "LLM PROVIDERS — ANY MODEL, ANY VENDOR", 0.35, 1.1, 12.6, 0.3,
        size=11, bold=True, color=BLUE)
    providers = [
        ("OpenAI", "GPT-4, GPT-4o, GPT-4 Turbo"),
        ("Anthropic", "Claude 3 Opus/Sonnet/Haiku, Claude 4"),
        ("Google", "Gemini 1.5 Pro/Flash"),
        ("Azure", "Azure OpenAI Service"),
        ("AWS Bedrock", "Meta Llama, Titan, Mistral"),
        ("Mistral", "Mistral Large/Medium"),
        ("Cohere", "Command R+"),
        ("Self-hosted", "Any open-source model"),
    ]
    for i, (name, models) in enumerate(providers):
        col_n = i % 4
        row_n = i // 4
        x = 0.35 + col_n * 3.22
        y = 1.48 + row_n * 0.88
        rect(s, x, y, 3.05, 0.78, CARD)
        rect(s, x, y, 3.05, 0.08, BLUE)
        txt(s, name, x + 0.1, y + 0.12, 2.85, 0.3, size=12, bold=True, color=WHITE)
        txt(s, models, x + 0.1, y + 0.42, 2.85, 0.3, size=9, color=GRAY)

    # Enterprise connectors
    txt(s, "ENTERPRISE CONNECTORS (50+)", 0.35, 3.32, 12.6, 0.3,
        size=11, bold=True, color=TEAL)
    connectors = [
        ("CRM",          "Salesforce · HubSpot · MS Dynamics",        TEAL),
        ("ERP",          "SAP · Oracle · Workday",                     TEAL),
        ("ITSM",         "ServiceNow · Jira · Zendesk",                TEAL),
        ("HRIS",         "Workday · BambooHR · ADP",                   TEAL),
        ("Legal",        "Westlaw · LexisNexis · CourtListener",       BLUE),
        ("Healthcare",   "Epic · Cerner · FDA APIs · PubMed",          BLUE),
        ("Financial",    "Bloomberg · Reuters · Core Banking",          BLUE),
        ("Identity",     "Active Directory · Okta · Azure AD",          BLUE),
    ]
    for i, (cat, systems, col) in enumerate(connectors):
        c_n = i % 4
        r_n = i // 4
        x = 0.35 + c_n * 3.22
        y = 3.7 + r_n * 0.82
        rect(s, x, y, 3.05, 0.72, CARD)
        rect(s, x, y, 0.05, 0.72, col)
        txt(s, cat, x + 0.15, y + 0.06, 2.8, 0.28, size=12, bold=True, color=col)
        txt(s, systems, x + 0.15, y + 0.38, 2.8, 0.28, size=9, color=GRAY)

    # Integration modes
    rect(s, 0.35, 5.45, 12.6, 1.75, CARD2)
    rect(s, 0.35, 5.45, 0.06, 1.75, BLUE)
    txt(s, "INTEGRATION MODES", 0.55, 5.52, 3, 0.3, size=11, bold=True, color=BLUE)
    modes = [
        ("SDK",         "hours",   "Import library; wrap LLM calls"),
        ("API Gateway", "minutes", "Change base URL — zero code changes"),
        ("On-Premise",  "1-2 wks", "Deploy in customer VPC"),
        ("SaaS",        "minutes", "Fully managed; instant start"),
    ]
    for i, (mode, time, desc) in enumerate(modes):
        x = 0.55 + i * 3.1
        txt(s, mode, x, 5.88, 3.0, 0.35, size=14, bold=True, color=WHITE)
        tag(s, x, 6.28, time, BLUE, w=1.3, h=0.27)
        txt(s, desc, x, 6.62, 3.0, 0.48, size=10, color=GRAY)

# ============================================================================
# SLIDE 17 — TECHNOLOGY STACK
# ============================================================================
def s17_tech_stack(prs):
    s = slide(prs)
    hdr(s, "Technology Stack",
        "Production-grade architecture built for enterprise scale")

    layers = [
        ("Application Layer", BLUE, [
            ("Demo Frontend",    "Streamlit + Plotly",             "Interactive dashboard & visualizations"),
            ("API Gateway",      "FastAPI (Python)",               "High-performance async request routing"),
            ("Detection Engine", "Python + HuggingFace Transformers","Parallel NLP-based detection pipeline"),
            ("Scoring Service",  "scikit-learn ensemble",          "Multi-technique result aggregation"),
        ]),
        ("Detection & ML", TEAL, [
            ("Semantic Analysis",    "sentence-transformers",      "Embedding-based entropy measurement"),
            ("Claim Extraction",     "spaCy NER + custom models",  "Named entity + assertion extraction"),
            ("Pattern Recognition",  "Fine-tuned RoBERTa",         "Hallucination signature classifier"),
            ("Numerical Validation", "Regex + NER pipeline",       "Numerical entity extraction & lookup"),
        ]),
        ("Data & Persistence", ORANGE, [
            ("Audit Logging",        "PostgreSQL (append-only)",   "Immutable interaction audit trail"),
            ("Real-time Metrics",    "TimescaleDB / InfluxDB",     "Time-series performance monitoring"),
            ("Enterprise Cache",     "Redis (TTL-based)",          "Fresh enterprise data snapshots"),
            ("Compliance Reports",   "Auto-generated PDF/JSON",    "Regulatory documentation output"),
        ]),
        ("Infrastructure", PURPLE, [
            ("Containerization",  "Docker + Docker Compose",       "Reproducible environments"),
            ("Orchestration",     "Kubernetes (production)",       "Auto-scaling, self-healing"),
            ("Cloud",             "AWS / Azure / GCP",             "Multi-cloud, customer-choice"),
            ("CI/CD",             "GitHub Actions",                "Automated test + deploy pipeline"),
        ]),
    ]

    for i, (layer, col, components) in enumerate(layers):
        col_n = i % 2
        row_n = i // 2
        x = 0.35 + col_n * 6.55
        y = 1.1 + row_n * 2.95

        rect(s, x, y, 6.25, 2.75, CARD)
        rect(s, x, y, 6.25, 0.45, col)
        txt(s, layer, x, y, 6.25, 0.45, size=14, bold=True, align=PP_ALIGN.CENTER)

        # Column headers
        for j, (hdr_t, hdr_x_off, hdr_w) in enumerate([
            ("Component", 0.1, 1.8),
            ("Technology", 1.95, 1.9),
            ("Purpose", 3.9, 2.25),
        ]):
            txt(s, hdr_t, x + hdr_x_off, y + 0.5, hdr_w, 0.25, size=9, bold=True, color=col)

        for k, (comp, tech, purpose) in enumerate(components):
            ry = y + 0.8 + k * 0.47
            row_bg = CARD2 if k % 2 == 0 else CARD
            rect(s, x + 0.1, ry, 6.05, 0.42, row_bg)
            txt(s, comp,    x + 0.15, ry + 0.07, 1.75, 0.28, size=9.5, color=WHITE)
            txt(s, tech,    x + 2.05, ry + 0.07, 1.85, 0.28, size=9.5, color=col)
            txt(s, purpose, x + 3.95, ry + 0.07, 2.2,  0.28, size=9,   color=GRAY)

# ============================================================================
# SLIDE 18 — ROI & BUSINESS CASE
# ============================================================================
def s18_roi(prs):
    s = slide(prs)
    hdr(s, "ROI & Business Case",
        "The math is simple — one prevented incident pays for TrustLayer AI many times over")

    # Cost model
    txt(s, "BASELINE (100,000 AI queries/month, no TrustLayer)", 0.35, 1.1, 12.6, 0.32,
        size=11, bold=True, color=RED)
    baseline = [
        ("Monthly AI query volume",     "100,000"),
        ("Industry hallucination rate", "15%  =  15,000 hallucinated responses/month"),
        ("Critical hallucinations",     "~25%  =  ~3,750 potentially harmful responses/month"),
        ("Critical incidents/year",     "~30  (incidents that cause measurable harm)"),
        ("Average cost per incident",   "$2,400,000  (remediation + legal + regulatory + reputation)"),
        ("TOTAL ANNUAL EXPOSURE",       "~$72,000,000"),
    ]
    for i, (label, val) in enumerate(baseline):
        y = 1.48 + i * 0.38
        row_bg = CARD if i % 2 == 0 else CARD2
        rect(s, 0.35, y, 12.6, 0.35, row_bg)
        if i == 5:
            rect(s, 0.35, y, 12.6, 0.35, RGBColor(0x2A, 0x10, 0x10))
            rect(s, 0.35, y, 0.05, 0.35, RED)
        txt(s, label, 0.52, y + 0.05, 5.5, 0.26, size=11, color=WHITE if i < 5 else RED, bold=(i==5))
        txt(s, val, 6.1, y + 0.05, 6.7, 0.26, size=11, color=LGRAY if i < 5 else RED, bold=(i==5))

    # With TrustLayer
    txt(s, "WITH TRUSTLAYER AI", 0.35, 3.82, 12.6, 0.32, size=11, bold=True, color=TEAL)
    with_tl = [
        ("Hallucination catch rate",  "92%  →  ~28 of 30 critical incidents prevented"),
        ("Incidents remaining/year",  "~2  (8% that pass through)"),
        ("Annual losses prevented",   "~$66,400,000"),
        ("TrustLayer annual cost",    "~$500,000  (enterprise tier SaaS)"),
        ("NET ANNUAL BENEFIT",        "~$65,900,000"),
        ("ROI",                       "> 400%  in year 1"),
    ]
    for i, (label, val) in enumerate(with_tl):
        y = 4.2 + i * 0.38
        row_bg = CARD if i % 2 == 0 else CARD2
        rect(s, 0.35, y, 12.6, 0.35, row_bg)
        is_highlight = i >= 4
        if is_highlight:
            rect(s, 0.35, y, 12.6, 0.35, RGBColor(0x00, 0x2A, 0x1A))
            rect(s, 0.35, y, 0.05, 0.35, TEAL)
        txt(s, label, 0.52, y + 0.05, 5.5, 0.26, size=11, color=WHITE, bold=is_highlight)
        txt(s, val, 6.1, y + 0.05, 6.7, 0.26, size=11, color=TEAL if is_highlight else LGRAY, bold=is_highlight)

    # Compliance upside
    rect(s, 0.35, 6.5, 12.6, 0.7, RGBColor(0x00, 0x15, 0x40))
    rect(s, 0.35, 6.5, 0.06, 0.7, BLUE)
    txt(s, "Additional upside — EU AI Act compliance avoids fines up to €35M (August 2026 enforcement).  "
        "Reputation protection.  Board-level AI governance readiness.",
        0.55, 6.6, 12.1, 0.5, size=12, color=WHITE)

# ============================================================================
# SLIDE 19 — ALL 13 DEMO SCENARIOS
# ============================================================================
def s19_demo_scenarios(prs):
    s = slide(prs)
    hdr(s, "Live Demo — All 13 Scenarios", "Real hallucinations. Real detection. Real-time decisions.")

    # Two-column layout
    scenarios = [
        ("BFSI — Banking",       "Mortgage rate inquiry",                "18%", "BLOCK", RED),
        ("BFSI — Insurance",     "Claims processing",                    "22%", "BLOCK", RED),
        ("BFSI — Wealth Mgmt",   "Investment advice",                    "31%", "BLOCK", RED),
        ("Healthcare — Clinical","Drug dosage (metformin)",               "6%",  "BLOCK", RED),
        ("Healthcare — Admin",   "Prior authorization",                  "29%", "BLOCK", RED),
        ("Legal — Litigation",   "Case research (all citations fabricated)","4%", "BLOCK", RED),
        ("Legal — Contract",     "NDA review (no document provided)",    "35%", "BLOCK", RED),
        ("Enterprise — HR",      "Policy inquiry (vacation / overtime)", "42%", "FLAG",  ORANGE),
        ("Enterprise — IT",      "VPN support (fabricated URLs)",        "38%", "FLAG",  ORANGE),
        ("Manufacturing — QC",   "Quality control record query",         "55%", "FLAG",  ORANGE),
        ("Retail — Product",     "Pricing inquiry (unverified discount)", "61%", "FLAG",  ORANGE),
        ("Telecom — Billing",    "Account modification query",           "27%", "BLOCK", RED),
        ("Government — Benefits","Benefits eligibility determination",   "19%", "BLOCK", RED),
    ]

    # Headers
    for col_n in range(2):
        x_off = col_n * 6.65
        hdrs = ["Industry", "Scenario", "Confidence", "Action"]
        hx = [0.35, 2.05, 4.85, 6.0]
        hw = [1.65, 2.75, 1.1, 0.95]
        for j, (h, hx_v, hw_v) in enumerate(zip(hdrs, hx, hw)):
            rect(s, x_off + hx_v, 1.1, hw_v, 0.38, BLUE)
            txt(s, h, x_off + hx_v + 0.05, 1.1, hw_v - 0.1, 0.38,
                size=10, bold=True, align=PP_ALIGN.CENTER)

    rows_per_col = 7
    for i, (industry, scenario, conf, action, col) in enumerate(scenarios):
        col_n = i // rows_per_col
        row_n = i % rows_per_col
        x_off = col_n * 6.65
        y = 1.52 + row_n * 0.76
        row_bg = CARD if row_n % 2 == 0 else CARD2

        hx = [0.35, 2.05, 4.85, 6.0]
        hw = [1.65, 2.75, 1.1, 0.95]
        for hx_v, hw_v in zip(hx, hw):
            rect(s, x_off + hx_v, y, hw_v, 0.68, row_bg)

        txt(s, industry,  x_off + hx[0] + 0.08, y + 0.1, hx[1] - 0.1, 0.48, size=10, color=col)
        txt(s, scenario,  x_off + hx[1] + 0.08, y + 0.1, hx[2] - 0.1, 0.48, size=10, color=WHITE)
        txt(s, conf,      x_off + hx[2] + 0.05, y + 0.12, hw[2] - 0.1, 0.42,
            size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, action,    x_off + hx[3] + 0.02, y + 0.12, hw[3] - 0.04, 0.42,
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 20 — DEMO FLOW & PLATFORM PAGES
# ============================================================================
def s20_demo_flow(prs):
    s = slide(prs)
    hdr(s, "Demo Application — 9 Pages, 15-Minute Flow",
        "Streamlit-based interactive demo covering all platform capabilities")

    pages = [
        ("01", "Executive Overview",      BLUE,   "Market stats · Capability cards · Industry risk heatmap"),
        ("02", "Live Detection Demo",     RED,    "13 interactive scenarios · Real-time scoring animation"),
        ("03", "Technical Architecture",  ORANGE, "8 detection techniques · Latency metrics · System diagram"),
        ("04", "Industry Solutions",      TEAL,   "8 vertical modules · Use cases · Integration details"),
        ("05", "Real-Time Monitoring",    GREEN,  "Live query volume · Detection feed · Alert stream"),
        ("06", "Compliance & Governance", PURPLE, "EU AI Act · NIST · ISO · SOC 2 coverage meters"),
        ("07", "Integrations",            BLUE,   "8 LLM providers · 50+ enterprise connector catalog"),
        ("08", "ROI Calculator",          TEAL,   "Customizable inputs · 12-month projection · Break-even"),
        ("09", "API Documentation",       ORANGE, "SDK examples · API reference · Integration patterns"),
    ]

    for i, (num, title, col, desc) in enumerate(pages):
        col_n = i % 3
        row_n = i // 3
        x = 0.35 + col_n * 4.35
        y = 1.1 + row_n * 1.6
        rect(s, x, y, 4.15, 1.42, CARD)
        rect(s, x, y, 4.15, 0.08, col)
        txt(s, num, x + 0.12, y + 0.12, 0.6, 0.52, size=22, bold=True, color=col)
        txt(s, title, x + 0.82, y + 0.12, 3.2, 0.4, size=12, bold=True, color=WHITE)
        txt(s, desc, x + 0.12, y + 0.72, 3.9, 0.62, size=10, color=GRAY)

    # Demo flow timeline
    txt(s, "RECOMMENDED 15-MINUTE DEMO FLOW", 0.35, 6.0, 12.6, 0.3, size=11, bold=True, color=WHITE)
    flow = [
        ("0-2m",   "Executive\nOverview",  BLUE),
        ("2-6m",   "Live Detection\nDemo", RED),
        ("6-8m",   "Technical\nArchitecture", ORANGE),
        ("8-10m",  "Industry\nSolutions",  TEAL),
        ("10-12m", "Compliance\n& Governance", PURPLE),
        ("12-15m", "ROI\nCalculator",      GREEN),
    ]
    for i, (time, label, col) in enumerate(flow):
        x = 0.35 + i * 2.15
        rect(s, x, 6.35, 2.0, 0.95, CARD)
        rect(s, x, 6.35, 2.0, 0.08, col)
        txt(s, time, x + 0.08, 6.45, 1.85, 0.28, size=11, bold=True, color=col)
        txt(s, label, x + 0.08, 6.75, 1.85, 0.5, size=10, color=WHITE)

# ============================================================================
# SLIDE 21 — COMPETITIVE LANDSCAPE
# ============================================================================
def s21_competitive(prs):
    s = slide(prs)
    hdr(s, "Competitive Landscape",
        "TrustLayer AI is the only platform combining detection + prevention + governance")

    # Feature matrix
    features = ["Real-time detection", "Automated blocking", "Enterprise grounding",
                "Compliance (EU AI Act)", "8 industry modules", "LLM-agnostic",
                "Audit trail", "Time to value"]
    competitors = [
        ("OpenAI\nSafety",     ["-",  "-",  "-",  "-",  "-",  "-",  "~",  "Fast"],  RED),
        ("Guardrails\nAI",     ["✓",  "~",  "-",  "-",  "-",  "✓",  "-",  "Med"],   ORANGE),
        ("Arthur\nAI",         ["~",  "-",  "-",  "~",  "-",  "✓",  "✓",  "Slow"],  ORANGE),
        ("Fiddler\nAI",        ["~",  "-",  "-",  "~",  "-",  "✓",  "✓",  "Slow"],  ORANGE),
        ("Custom\nBuild",      ["✓",  "✓",  "~",  "~",  "-",  "✓",  "~",  "V.Slow"],RED),
        ("TrustLayer\nAI",     ["✓",  "✓",  "✓",  "✓",  "✓",  "✓",  "✓",  "Fast"],  TEAL),
    ]

    # Header column
    rect(s, 0.35, 1.1, 2.0, 0.42, DARK_BLUE)
    txt(s, "Feature", 0.45, 1.1, 1.9, 0.42, size=10, bold=True, align=PP_ALIGN.CENTER)
    for i, feat in enumerate(features):
        y = 1.57 + i * 0.6
        rect(s, 0.35, y, 2.0, 0.55, CARD if i % 2 == 0 else CARD2)
        txt(s, feat, 0.45, y + 0.1, 1.85, 0.38, size=10, color=WHITE)

    # Competitor columns
    col_w = 1.82
    for j, (comp, vals, col) in enumerate(competitors):
        x = 2.4 + j * col_w
        is_tl = j == 5
        hdr_bg = TEAL if is_tl else DARK_BLUE
        rect(s, x, 1.1, col_w - 0.05, 0.42, hdr_bg)
        txt(s, comp, x, 1.1, col_w - 0.05, 0.42, size=10, bold=True, align=PP_ALIGN.CENTER)
        for i, val in enumerate(vals):
            y = 1.57 + i * 0.6
            row_bg = RGBColor(0x00, 0x2A, 0x1A) if is_tl else (CARD if i % 2 == 0 else CARD2)
            rect(s, x, y, col_w - 0.05, 0.55, row_bg)
            val_col = TEAL if val == "✓" else (RED if val == "-" else ORANGE)
            if is_tl: val_col = TEAL
            txt(s, val, x, y + 0.1, col_w - 0.05, 0.38, size=12,
                bold=True, color=val_col, align=PP_ALIGN.CENTER)

    # Bottom summary
    rect(s, 0.35, 6.58, 12.6, 0.68, RGBColor(0x00, 0x22, 0x55))
    rect(s, 0.35, 6.58, 0.06, 0.68, TEAL)
    txt(s, "TrustLayer AI is the only solution delivering full coverage across all 8 dimensions — "
        "detection, prevention, governance, compliance, industry specialization, and rapid time-to-value. "
        "No other platform comes close.",
        0.55, 6.68, 12.1, 0.5, size=12, color=WHITE)

# ============================================================================
# SLIDE 22 — ROADMAP
# ============================================================================
def s22_roadmap(prs):
    s = slide(prs)
    hdr(s, "Product Roadmap",
        "From hackathon demo to enterprise-grade platform and beyond")

    phases = [
        ("Phase 1", "NOW", "Hackathon Demo", BLUE, [
            "Interactive Streamlit demo — 9 pages",
            "13 industry hallucination scenarios",
            "8 simulated detection techniques",
            "Compliance dashboards (EU AI Act, NIST, ISO, SOC 2)",
            "ROI calculator with customizable inputs",
        ]),
        ("Phase 2", "3–6 Mo", "MVP", TEAL, [
            "Production NLP detection engine",
            "Live API gateway (zero code integration)",
            "3 enterprise connectors (Salesforce, ServiceNow, SAP)",
            "Real audit logging + compliance report generation",
            "First 1–2 paying enterprise pilots",
        ]),
        ("Phase 3", "6–12 Mo", "Growth", ORANGE, [
            "50+ enterprise system connectors",
            "Industry-specific fine-tuned models",
            "SOC 2 Type II certification",
            "Multi-region (US / EU / APAC)",
            "Self-serve onboarding portal",
        ]),
        ("Phase 4", "12–24 Mo", "Scale", PURPLE, [
            "EU AI Act compliance certification program",
            "Third-party detection module marketplace",
            "Federated learning (privacy-preserving)",
            "Real-time regulatory update pipeline",
            "Series B fundraising milestone",
        ]),
    ]

    for i, (phase, when, title, col, items) in enumerate(phases):
        x = 0.35 + i * 3.25
        rect(s, x, 1.1, 3.1, 6.0, CARD)
        rect(s, x, 1.1, 3.1, 0.55, col)
        txt(s, phase, x + 0.12, 1.13, 1.8, 0.28, size=14, bold=True)
        tag(s, x + 2.05, 1.18, when, DARK_BLUE, w=0.9, h=0.28)
        txt(s, title, x + 0.12, 1.62, 2.86, 0.38, size=14, bold=True, color=col)
        for j, item in enumerate(items):
            y = 2.1 + j * 0.92
            rect(s, x + 0.12, y, 2.86, 0.78, CARD2)
            rect(s, x + 0.12, y, 0.04, 0.78, col)
            txt(s, item, x + 0.25, y + 0.12, 2.65, 0.55, size=10.5, color=WHITE)

# ============================================================================
# SLIDE 23 — WHY NOW / MARKET TIMING
# ============================================================================
def s23_why_now(prs):
    s = slide(prs)
    hdr(s, "Why Now — Perfect Market Timing",
        "Four converging forces make 2025-2026 the critical window")

    forces = [
        ("Regulatory Deadline",
         "EU AI Act Full Enforcement: August 2026",
         "Enterprises have 12-18 months to build compliance infrastructure.\n"
         "TrustLayer provides 94% coverage out of the box — immediate value.",
         RED),
        ("Enterprise AI Adoption Surge",
         "AI spending doubled in 2024; projected to triple by 2026",
         "More AI deployments = more hallucination risk = larger addressable market.\n"
         "The problem grows faster than awareness of the solution.",
         BLUE),
        ("High-Profile Incidents",
         "Legal sanctions, healthcare near-misses driving board attention",
         "Risk of AI hallucinations is now a board-level concern at 77% of enterprises.\n"
         "Boards are demanding answers — TrustLayer provides them.",
         ORANGE),
        ("Provider Inaction",
         "OpenAI, Google, Microsoft cannot solve this conflict-of-interest problem",
         "The window for an independent trust layer solution is wide open.\n"
         "First-mover advantage in a category that will be essential by 2026.",
         TEAL),
    ]

    for i, (title, headline, body, col) in enumerate(forces):
        x = 0.35 + (i % 2) * 6.5
        y = 1.1 + (i // 2) * 3.05
        rect(s, x, y, 6.2, 2.8, CARD)
        rect(s, x, y, 6.2, 0.55, col)
        txt(s, title, x + 0.15, y + 0.08, 5.9, 0.42, size=15, bold=True)
        txt(s, headline, x + 0.15, y + 0.65, 5.9, 0.35, size=12, bold=True, color=col)
        txt(s, body, x + 0.15, y + 1.1, 5.9, 1.55, size=11, color=LGRAY)

    # Market size
    rect(s, 0.35, 7.2, 12.6, 0.3, DARK_BLUE)
    txt(s, "AI Governance & Trust Market:  $2.1B (2024)  →  $18.7B projected (2030)  |  CAGR: 43%",
        0.5, 7.24, 12.1, 0.22, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 24 — CLOSING / CALL TO ACTION
# ============================================================================
def s24_closing(prs):
    s = slide(prs)
    rect(s, 0, 0, 13.33, 7.5, BG)
    rect(s, 0, 0, 0.12, 7.5, BLUE)
    rect(s, 0.12, 0, 13.21, 0.06, BLUE)

    txt(s, "TrustLayer AI", 0.5, 1.0, 12.5, 1.2, size=52, bold=True, align=PP_ALIGN.CENTER)
    txt(s, "Enterprise AI Reliability Platform",
        0.5, 2.15, 12.5, 0.55, size=20, color=LBLUE, align=PP_ALIGN.CENTER)
    rect(s, 2.5, 2.9, 8.3, 0.05, BLUE)

    txt(s,
        "Every day you deploy AI without reliability controls is a day you are exposed.\n"
        "Exposed to customer harm. Regulatory fines. Reputation damage.",
        0.5, 3.1, 12.5, 0.9, size=14, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Three pillars
    pillars = [
        ("DETECT",  BLUE,   "8 parallel algorithms\n92% catch rate\n<100ms latency"),
        ("PREVENT", TEAL,   "Real-time blocking\nbefore harm occurs\nzero end-user delay"),
        ("GOVERN",  PURPLE, "EU AI Act ready\nFull audit trail\nBoard-ready reports"),
    ]
    for i, (title, col, desc) in enumerate(pillars):
        x = 1.1 + i * 3.75
        rect(s, x, 4.2, 3.5, 2.1, CARD)
        rect(s, x, 4.2, 3.5, 0.55, col)
        txt(s, title, x, 4.2, 3.5, 0.55, size=20, bold=True, align=PP_ALIGN.CENTER)
        txt(s, desc, x + 0.15, 4.85, 3.2, 1.2, size=13, color=GRAY, align=PP_ALIGN.CENTER)

    rect(s, 0, 6.85, 13.33, 0.65, BLUE)
    txt(s,
        "TrustLayer AI  ·  Built for the industries that cannot afford to get it wrong"
        "  ·  Infosys Business Incubator Cohort 2",
        0.5, 6.95, 12.5, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================================
# MAIN
# ============================================================================

def build():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides = [
        ("01", "Title / Hero",                    s01_title),
        ("02", "Crisis Scale",                     s02_crisis_scale),
        ("03", "Real-World Incidents",             s03_incidents),
        ("04", "Root Cause & Market Gap",          s04_root_cause),
        ("05", "Solution Overview",                s05_solution),
        ("06", "System Architecture",              s06_architecture),
        ("07", "Detection Engine Overview",        s07_detection_overview),
        ("08", "Detection Deep-Dive",              s08_detection_deepdive),
        ("09", "Scoring Model & Decision Logic",   s09_scoring),
        ("10", "Industry Modules Overview",        s10_industries),
        ("11", "BFSI Banking Walkthrough",         s11_bfsi),
        ("12", "Healthcare Deep-Dive",             s12_healthcare),
        ("13", "Legal Deep-Dive",                  s13_legal),
        ("14", "Compliance & Governance",          s14_compliance),
        ("15", "EU AI Act Deep-Dive",              s15_eu_ai_act),
        ("16", "Integration Ecosystem",            s16_integrations),
        ("17", "Technology Stack",                 s17_tech_stack),
        ("18", "ROI & Business Case",              s18_roi),
        ("19", "All 13 Demo Scenarios",            s19_demo_scenarios),
        ("20", "Demo Application Flow",            s20_demo_flow),
        ("21", "Competitive Landscape",            s21_competitive),
        ("22", "Roadmap",                          s22_roadmap),
        ("23", "Why Now / Market Timing",          s23_why_now),
        ("24", "Closing / Call to Action",         s24_closing),
    ]

    print(f"Building {len(slides)}-slide detailed presentation...\n")
    for num, title, fn in slides:
        fn(prs)
        print(f"  Slide {num}  {title}")

    out = r"C:\aiprojects\aitrust\TrustLayer_AI_Hackathon_Detailed.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")


if __name__ == "__main__":
    build()
